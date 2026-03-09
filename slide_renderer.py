"""
slide_renderer.py — Renderizza slide PPTX come immagini PNG
============================================================

Produce un PNG per ogni slide del file .pptx.
Le immagini vengono salvate in images_dir con naming:
    {stem}_slide_001.png, {stem}_slide_002.png, ...

Il prefisso {stem} (nome del file PPTX senza estensione) evita collisioni
quando una lezione contiene più file .pptx.

Queste immagini hanno DUE usi distinti:
  1. Nel LaTeX finale — ogni \subsection include la slide come \begin{figure}
  2. Nel prompt a Claude — Claude vede visivamente la slide

Strategia di rendering (in ordine di priorità):
  1. python-pptx + Pillow  — rendering nativo, niente animazioni, robusto
  2. pymupdf (fitz)        — se il PPTX è già stato convertito in PDF
  3. placeholder PNG       — fallback se nessuna libreria disponibile

Uso:
    from slide_renderer import render_slide_images

    slide_images = render_slide_images(
        pptx_path  = Path("lezione_01.pptx"),
        images_dir = Path("output/images"),
    )
    # slide_images = {1: "lezione_slide_001.png", 2: "lezione_slide_002.png", ...}
"""

import os
import hashlib
from pathlib import Path
from typing import Optional


# ─────────────────────────────────────────────
# DIPENDENZE OPZIONALI
# ─────────────────────────────────────────────

def _check_deps() -> dict:
    """Verifica quali dipendenze sono disponibili."""
    deps = {"pptx": False, "pillow": False, "pymupdf": False}
    try:
        from pptx import Presentation
        deps["pptx"] = True
    except ImportError:
        pass
    try:
        from PIL import Image, ImageDraw
        deps["pillow"] = True
    except ImportError:
        pass
    try:
        import fitz
        deps["pymupdf"] = True
    except ImportError:
        pass
    return deps


# ─────────────────────────────────────────────
# RENDERING PRINCIPALE
# ─────────────────────────────────────────────

def render_slide_images(pptx_path: Path, images_dir: Path) -> dict:
    """
    Renderizza ogni slide come PNG.

    Ritorna dizionario:
        {slide_number: "{stem}_slide_001.png", ...}
    I path sono relativi a images_dir (pronti per \includegraphics).

    Se il rendering fallisce su una slide, quella slide viene saltata
    senza bloccare le altre.
    """
    images_dir = Path(images_dir)
    images_dir.mkdir(parents=True, exist_ok=True)

    deps = _check_deps()

    if deps["pptx"] and deps["pillow"]:
        print(f"    [renderer] python-pptx + Pillow")
        return _render_with_pptx_pillow(pptx_path, images_dir)
    elif deps["pymupdf"]:
        print(f"    [renderer] pymupdf (fallback)")
        return _render_with_pymupdf(pptx_path, images_dir)
    else:
        print(f"    [renderer] WARN: nessuna libreria disponibile")
        print(f"               pip install Pillow  oppure  pip install pymupdf")
        return _render_placeholder(pptx_path, images_dir)


# ─────────────────────────────────────────────
# METODO 1: python-pptx + Pillow
# Disegna ogni shape sulla slide elemento per elemento.
# Ignora le animazioni → mostra tutto il contenuto nello stato finale.
# ─────────────────────────────────────────────

def _render_with_pptx_pillow(pptx_path: Path, images_dir: Path) -> dict:
    from pptx import Presentation
    from pptx.util import Emu
    from PIL import Image, ImageDraw, ImageFont
    import io

    prs    = Presentation(str(pptx_path))
    stem   = Path(pptx_path).stem
    result = {}

    # Dimensioni slide in pixel (96 DPI standard)
    DPI    = 150
    EMU_PER_INCH = 914400
    slide_w_px = int(prs.slide_width  / EMU_PER_INCH * DPI)
    slide_h_px = int(prs.slide_height / EMU_PER_INCH * DPI)

    def emu_to_px(emu: int) -> int:
        return int(emu / EMU_PER_INCH * DPI)

    for slide_idx, slide in enumerate(prs.slides, start=1):
        img_filename = f"{stem}_slide_{slide_idx:03d}.png"
        img_path     = images_dir / img_filename

        # Usa cache se il file esiste già
        if img_path.exists():
            result[slide_idx] = img_filename
            continue

        try:
            # Canvas bianco
            canvas = Image.new("RGB", (slide_w_px, slide_h_px), "white")
            draw   = ImageDraw.Draw(canvas)

            # Sfondo slide (se ha un colore di riempimento)
            try:
                bg = slide.background
                fill = bg.fill
                if fill.type is not None:
                    try:
                        from pptx.dml.color import RGBColor
                        rgb = fill.fore_color.rgb
                        canvas = Image.new(
                            "RGB", (slide_w_px, slide_h_px),
                            (rgb[0], rgb[1], rgb[2])
                        )
                        draw = ImageDraw.Draw(canvas)
                    except Exception:
                        pass
            except Exception:
                pass

            # Disegna ogni shape
            for shape in slide.shapes:
                try:
                    left   = emu_to_px(shape.left   or 0)
                    top    = emu_to_px(shape.top    or 0)
                    width  = emu_to_px(shape.width  or 0)
                    height = emu_to_px(shape.height or 0)

                    # Immagine embedded
                    if shape.shape_type == 13:
                        try:
                            img_bytes = shape.image.blob
                            embedded  = Image.open(io.BytesIO(img_bytes)).convert("RGBA")
                            # Preserva aspect ratio: scala dentro il bounding box
                            if width > 0 and height > 0:
                                embedded.thumbnail((width, height), Image.LANCZOS)
                                # Centra nell'area della shape
                                paste_x = left + (width  - embedded.width)  // 2
                                paste_y = top  + (height - embedded.height) // 2
                            else:
                                paste_x, paste_y = left, top
                            canvas.paste(embedded, (paste_x, paste_y), embedded)
                        except Exception:
                            # Box grigio placeholder per immagine non renderizzabile
                            draw.rectangle(
                                [left, top, left+width, top+height],
                                fill="#e0e0e0", outline="#999999"
                            )
                        continue

                    # Testo
                    if shape.has_text_frame:
                        y_offset = top + 4
                        for para in shape.text_frame.paragraphs:
                            line = "".join(r.text for r in para.runs).strip()
                            if not line:
                                y_offset += 8
                                continue

                            # Dimensione font: leggi dai run pptx, scala a DPI
                            pt_size = 12
                            try:
                                for run in para.runs:
                                    if run.font.size:
                                        pt_size = int(run.font.size.pt)
                                        break
                                if para.runs and para.runs[0].font.size is None:
                                    # Prova dal placeholder / theme
                                    if shape.text_frame.paragraphs[0].runs:
                                        pass  # usa il default
                            except Exception:
                                pass
                            font_px = max(8, int(pt_size * DPI / 72))

                            # Colore testo: leggi dal primo run con colore esplicito
                            text_color = "black"
                            try:
                                for run in para.runs:
                                    if run.font.color and run.font.color.type is not None:
                                        rgb = run.font.color.rgb
                                        text_color = (rgb[0], rgb[1], rgb[2])
                                        break
                            except Exception:
                                pass

                            # Carica font scalato
                            try:
                                font = ImageFont.load_default(size=font_px)
                            except TypeError:
                                # Pillow < 10 non supporta size=
                                font = ImageFont.load_default()

                            # Word-wrap: spezza la riga per stare nella box
                            box_w = width - 8 if width > 8 else slide_w_px - 8
                            words  = line.split()
                            wrapped_lines = []
                            current = ""
                            for word in words:
                                test = (current + " " + word).strip()
                                try:
                                    tw = draw.textlength(test, font=font)
                                except AttributeError:
                                    tw = len(test) * font_px * 0.6  # stima
                                if tw <= box_w or not current:
                                    current = test
                                else:
                                    wrapped_lines.append(current)
                                    current = word
                            if current:
                                wrapped_lines.append(current)

                            line_h = font_px + 3
                            for wl in wrapped_lines:
                                if height > 0 and y_offset > top + height - line_h:
                                    break
                                draw.text((left + 4, y_offset), wl,
                                          fill=text_color, font=font)
                                y_offset += line_h
                            y_offset += 2  # piccola spaziatura tra paragrafi

                except Exception:
                    pass  # Singola shape fallita → continua con le altre

            canvas.save(str(img_path), "PNG", optimize=True)
            result[slide_idx] = img_filename
            print(f"    ✓ {img_filename}  ({slide_w_px}×{slide_h_px}px)")

        except Exception as e:
            print(f"    [WARN] slide {slide_idx} non renderizzata: {e}")

    print(f"    ✓ Renderizzate {len(result)}/{len(prs.slides)} slide")
    return result


# ─────────────────────────────────────────────
# METODO 2: pymupdf
# Usato se il PPTX è già disponibile come PDF.
# ─────────────────────────────────────────────

def _render_with_pymupdf(pptx_path: Path, images_dir: Path) -> dict:
    """
    Converte prima in PDF con LibreOffice, poi renderizza con pymupdf.
    """
    import subprocess
    import fitz  # pymupdf

    result   = {}
    stem     = Path(pptx_path).stem
    pdf_path = images_dir / (stem + "_slides.pdf")

    # Converti PPTX → PDF
    if not pdf_path.exists():
        try:
            subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pdf",
                 "--outdir", str(images_dir), str(pptx_path)],
                capture_output=True, timeout=60
            )
        except Exception as e:
            print(f"    [ERRORE] LibreOffice: {e}")
            return {}

    if not pdf_path.exists():
        print(f"    [ERRORE] PDF non generato da LibreOffice")
        return {}

    # Renderizza ogni pagina del PDF
    doc = fitz.open(str(pdf_path))
    try:
        for page_idx, page in enumerate(doc, start=1):
            img_filename = f"{stem}_slide_{page_idx:03d}.png"
            img_path     = images_dir / img_filename
            if img_path.exists():
                result[page_idx] = img_filename
                continue
            mat = fitz.Matrix(150/72, 150/72)  # 150 DPI
            pix = page.get_pixmap(matrix=mat)
            pix.save(str(img_path))
            result[page_idx] = img_filename
            print(f"    ✓ {img_filename}")
    finally:
        doc.close()
    return result


# ─────────────────────────────────────────────
# FALLBACK: placeholder PNG
# Genera PNG grigi con numero slide — meglio di niente.
# ─────────────────────────────────────────────

def _render_placeholder(pptx_path: Path, images_dir: Path) -> dict:
    """
    Genera PNG placeholder (rettangolo grigio con numero slide).
    Usa pptx per il conteggio slide se disponibile; fallback a 1 placeholder.
    """
    try:
        from pptx import Presentation
        prs    = Presentation(str(pptx_path))
        slides = list(prs.slides)
    except ImportError:
        slides = [None]  # pptx non disponibile: genera 1 placeholder

    stem   = Path(pptx_path).stem
    result = {}

    for slide_idx, slide in enumerate(slides, start=1):
        img_filename = f"{stem}_slide_{slide_idx:03d}.png"
        img_path     = images_dir / img_filename
        if img_path.exists():
            result[slide_idx] = img_filename
            continue

        try:
            from PIL import Image, ImageDraw, ImageFont
            img  = Image.new("RGB", (960, 540), "#f5f5f5")
            draw = ImageDraw.Draw(img)
            draw.rectangle([0, 0, 959, 539], outline="#cccccc", width=2)

            title = ""
            try:
                title = (slide.shapes.title.text or "") if slide.shapes.title else ""
            except Exception:
                pass

            draw.text((480, 260), f"Slide {slide_idx}", fill="#999999", anchor="mm")
            if title:
                draw.text((480, 290), title[:60], fill="#666666", anchor="mm")

            img.save(str(img_path), "PNG")
            result[slide_idx] = img_filename

        except ImportError:
            # Crea file PNG minimo valido (1x1 pixel bianco)
            img_path.write_bytes(
                b'\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01'
                b'\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde\x00\x00'
                b'\x00\x0cIDATx\x9cc\xf8\x0f\x00\x00\x01\x01\x00\x05\x18'
                b'\xd8N\x00\x00\x00\x00IEND\xaeB`\x82'
            )
            result[slide_idx] = img_filename

    return result


# ─────────────────────────────────────────────
# HELPER — LaTeX figure block per una slide
# ─────────────────────────────────────────────

def slide_figure_latex(img_filename: str, slide_number: int,
                        caption: str = "") -> str:
    """
    Genera il blocco LaTeX \begin{figure} per una slide.
    Da inserire all'inizio di ogni \subsection.

    Esempio output:
        \\begin{figure}[H]
          \\centering
          \\includegraphics[width=\\textwidth]{slide_003.png}
          \\caption{Slide 3: Titolo}
        \\end{figure}
    """
    cap = f"  \\caption{{Slide {slide_number}" + (f": {caption}" if caption else "") + "}\n"
    return (
        f"\\begin{{figure}}[H]\n"
        f"  \\centering\n"
        f"  \\includegraphics[width=\\textwidth]{{{img_filename}}}\n"
        f"{cap}"
        f"\\end{{figure}}\n"
    )