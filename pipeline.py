#!/usr/bin/env python3
"""
pipeline.py — Orchestratore Appunti AI
=======================================
Converte fonti miste in un libro LaTeX strutturato.

Fonti supportate:
  Audio:  .mp3 .wav .m4a .ogg .flac
  Video:  .mp4 .mkv .avi .mov .webm  -> ffmpeg -> Whisper (auto-detect lingua)
  Slide:  .pptx                       -> extractor + omml2latex + pix2tex
  Word:   .docx                       -> python-docx (paragrafi + tabelle)
  PDF:    .pdf                        -> pdfplumber
  Testo:  .txt .md .rtf               -> diretto (RTF auto-stripped)

Output:
  output/
  ├── main.tex          <- compila questo con pdflatex
  ├── lezione_01.tex
  ├── lezione_02.tex
  └── images/

Uso:
  python pipeline.py ./lezione_01/
  python pipeline.py --batch ./corso/ --title "Analisi Matematica 1"
  python pipeline.py --batch ./corso/ --skip-ai --skip-ocr   # veloce, offline
"""

# definisce percorsi, estensioni dei file supportati, modelli AI da usare.
import argparse
import os
import sys
import json
import time
import subprocess
import shutil
from pathlib import Path
from datetime import datetime

# Aggiungi la cartella dello script al path per trovare i moduli del collega
SCRIPT_DIR = Path(__file__).parent.resolve()
sys.path.insert(0, str(SCRIPT_DIR))

# ─────────────────────────────────────────────
# CONFIGURAZIONE
# ─────────────────────────────────────────────
CONFIG = {
    "whisper_model":     "base",
    "claude_model":      "claude-sonnet-4-6",
    "claude_max_tokens": 32000,
    "ext_audio":  [".mp3", ".wav", ".m4a", ".ogg", ".flac"],
    "ext_video":  [".mp4", ".mkv", ".avi", ".mov", ".webm"],
    "ext_slide":  [".pptx"],
    "ext_doc":    [".docx"],
    "ext_pdf":    [".pdf"],
    "ext_text":   [".txt", ".md", ".rtf"],
    "images_subdir": "images",
}

# ─────────────────────────────────────────────
# IMPORT MODULI extractor builder formula_detector omml2latex ocr_math
# ─────────────────────────────────────────────
try:
    from extractor import extract_slides
    from builder import _escape_latex
    from formula_detector import is_formula_image
    from omml2latex import omml_to_latex
    from ocr_math import image_to_latex, unload_models as _ocr_unload
    from slide_renderer import render_slide_images, slide_figure_latex
    from pdf_renderer import render_pdf_pages, _build_pdf_latex_skeleton
    COLLEAGUE_MODULES = True
    print("✓ Moduli collega: extractor, builder, formula_detector, omml2latex, ocr_math, slide_renderer, pdf_renderer")
except ImportError as e:
    COLLEAGUE_MODULES = False
    print(f"⚠  Moduli collega non disponibili ({e}) — uso fallback base")

    # Fallback: _escape_latex non importata da builder.py → definizione inline
    # Usata da process_pdf_chunked e altri punti che richiedono escape LaTeX
    import re as _re_esc
    def _escape_latex(t: str) -> str:  # noqa: F811
        _MAP = {"\\": "\\textbackslash{}", "&": "\\&", "%": "\\%", "$": "\\$",
                "#": "\\#", "_": "\\_", "^": "\\^{}", "{": "\\{", "}": "\\}",
                "~": "\\textasciitilde{}"}
        return _re_esc.compile(r'[\\&%$#_^{}~]').sub(lambda m: _MAP[m.group()], t)


# ───────────────────────────────────────────────────────────────────────────────────
#                       ESTRAZIONE TESTO DA AUDIO (ffmpeg + whisper)
# ───────────────────────────────────────────────────────────────────────────────────

# ─────────────────────────────────────────────
# STEP 1: ESTRAZIONE AUDIO DA VIDEO (ffmpeg)
# ─────────────────────────────────────────────
def extract_audio_from_video(video_path: Path, out_dir: Path):
    mp3_path = out_dir / (video_path.stem + "_audio.mp3")
    if mp3_path.exists():
        print(f"    [cache] audio già estratto: {mp3_path.name}")
        return mp3_path
    print(f"    ffmpeg: {video_path.name} -> mp3 ...")
    cmd = ["ffmpeg", "-y", "-i", str(video_path),
           "-vn", "-ac", "1", "-codec:a", "libmp3lame", "-qscale:a", "4",
           str(mp3_path)]
    r = subprocess.run(cmd, capture_output=True)
    if r.returncode != 0:
        print(f"    [ERRORE] ffmpeg: {r.stderr.decode()[:200]}")
        return None
    print(f"    ✓ {mp3_path.name}")
    return mp3_path


# ─────────────────────────────────────────────
# STEP 2: TRASCRIZIONE WHISPER
# ─────────────────────────────────────────────
def _get_audio_duration(audio_path: Path) -> float:
    """Stima la durata audio in secondi tramite ffprobe (0 se non disponibile)."""
    try:
        r = subprocess.run(
            ["ffprobe", "-v", "error", "-show_entries", "format=duration",
             "-of", "csv=p=0", str(audio_path)],
            capture_output=True, text=True, timeout=10,
        )
        return float(r.stdout.strip() or 0)
    except Exception:
        return 0.0


def transcribe_audio(audio_path: Path, model_name: str = "base",
                     _prog_dir: Path = None, _prog_base: int = 10,
                     _prog_end: int = 30) -> tuple[str, int]:
    """
    Trascrive l'audio con Whisper.
    Ritorna (testo_con_timestamp, durata_secondi).
    _prog_dir / _prog_base / _prog_end: se forniti, scrive aggiornamenti periodici.
    """
    cache      = audio_path.with_suffix(".transcript.txt")
    cache_dur  = audio_path.with_suffix(".duration.txt")
    if cache.exists():
        print(f"    [cache] {cache.name}")
        duration = int(cache_dur.read_text().strip()) if cache_dur.exists() else 0
        return cache.read_text(encoding="utf-8"), duration
    try:
        import whisper
    except ImportError:
        print("    [MANCANTE] whisper — pip install openai-whisper")
        return "", 0

    # Stima durata per heartbeat
    audio_dur = _get_audio_duration(audio_path)
    # Whisper base: ~1× realtime; large: ~0.3×. Usiamo 1× come stima conservativa.
    estimated_sec = audio_dur if audio_dur > 0 else 600.0

    print(f"    Whisper ({model_name}): {audio_path.name} ...")
    t0 = time.time()

    # Thread heartbeat: aggiorna progress ogni 15s durante la trascrizione
    if _prog_dir:
        import threading
        _stop = threading.Event()
        def _heartbeat():
            while not _stop.wait(timeout=15):
                elapsed = time.time() - t0
                frac    = min(elapsed / estimated_sec, 0.95)
                pct     = int(_prog_base + frac * (_prog_end - _prog_base))
                _report_progress(_prog_dir, pct,
                                 f"Whisper — {audio_path.name}",
                                 f"{int(elapsed)}s / ~{int(estimated_sec)}s stimati")
        threading.Thread(target=_heartbeat, daemon=True).start()

    model = whisper.load_model(model_name)
    # Lingue supportate: en (default) e it. Qualsiasi altra lingua rilevata → forzata a en.
    # WHISPER_LANG=it per forzare italiano; senza variabile: auto-detect tra en/it.
    _ALLOWED_LANGS = {"en", "it"}
    whisper_lang = os.environ.get("WHISPER_LANG", "").strip().lower() or None
    if whisper_lang and whisper_lang not in _ALLOWED_LANGS:
        print(f"    [WARN] WHISPER_LANG='{whisper_lang}' non supportato → uso 'en'")
        whisper_lang = "en"
    # fp16 richiede CUDA; su CPU deve essere False altrimenti Whisper crasha
    try:
        import torch as _torch
        _fp16 = _torch.cuda.is_available()
    except ImportError:
        _fp16 = False
    result = model.transcribe(str(audio_path), language=whisper_lang, verbose=False, fp16=_fp16)
    # Valida lingua rilevata: se non en/it (es. turco), ri-trascrivi forzando en
    _detected = result.get("language", "en")
    if _detected not in _ALLOWED_LANGS:
        print(f"    [WARN] Lingua rilevata '{_detected}' non supportata → ri-trascrivo con 'en'")
        result = model.transcribe(str(audio_path), language="en", verbose=False, fp16=_fp16)

    if _prog_dir:
        _stop.set()
    elapsed = time.time() - t0
    lines = []
    duration = 0
    for seg in result.get("segments", []):
        m, s = int(seg["start"] // 60), int(seg["start"] % 60)
        lines.append(f"[{m:02d}:{s:02d}] {seg['text'].strip()}")
        duration = max(duration, int(seg.get("end", seg["start"])))
    text = "\n".join(lines)
    cache.write_text(text, encoding="utf-8")
    cache_dur.write_text(str(duration), encoding="utf-8")
    print(f"    ✓ {elapsed:.0f}s, {len(lines)} segmenti, durata≈{duration}s")
    return text, duration



# ───────────────────────────────────────────────────────────────────────────────────
#                                        PIPELINE PPTX / DOCX / PDF
# ───────────────────────────────────────────────────────────────────────────────────
def process_pptx_full(pptx_path: Path, images_dir: Path, skip_ocr: bool = False):
    """
    Usa: extractor -> omml2latex -> formula_detector + pix2tex (opz.)
    Ritorna (slides_list, testo_plain_per_claude)
    """
    print(f"    extractor: {pptx_path.name} ...")
    slides = extract_slides(str(pptx_path), str(images_dir))
    n_omml = n_ocr = 0

    # OMML -> LaTeX
    for slide in slides:
        for obj in slide.objects:
            if obj.obj_type == "omml_formula":
                obj.latex_result = omml_to_latex(obj.content)
                n_omml += 1

    # pix2tex: OCR su immagini che sembrano formule
    if not skip_ocr:
        candidates = [
            (s, o) for s in slides for o in s.objects
            if o.obj_type == "image" and o.image_path and is_formula_image(o.image_path)
        ]
        if candidates:
            print(f"    pix2tex: {len(candidates)} immagini formula ...")
            for _, obj in candidates:
                latex = image_to_latex(obj.image_path)
                if latex:
                    obj.latex_result = latex
                    n_ocr += 1
            # Libera modelli ML dopo il batch — evita memory leak su server long-running
            if COLLEAGUE_MODULES:
                _ocr_unload()

    total_obj = sum(len(s.objects) for s in slides)
    print(f"    ✓ {len(slides)} slide | {total_obj} oggetti | {n_omml} OMML | {n_ocr} OCR")

    # Testo plain per Claude
    lines = []
    n_notes = 0
    for slide in slides:
        lines.append(f"\n--- SLIDE {slide.slide_number}: {slide.title} ---")
        for obj in slide.objects:
            if obj.obj_type == "text" and obj.content.strip():
                lines.append(obj.content.strip())
            elif obj.obj_type == "omml_formula":
                f = getattr(obj, "latex_result", "")
                if f:
                    lines.append(f"[FORMULA: {f}]")
            elif obj.obj_type == "table" and obj.content.strip():
                # Invia la tabella LaTeX a Claude — la capisce e la riusa/migliora
                lines.append(f"[TABELLA:\n{obj.content}\n]")
        # Note del presentatore: fonte preziosa spesso ignorata
        if slide.notes:
            lines.append(f"[NOTE PRESENTER: {slide.notes}]")
            n_notes += 1

    if n_notes:
        print(f"    ✓ {n_notes} slide con note del presentatore estratte")
    return slides, "\n".join(lines)

# ─────────────────────────────────────────────
# STEP 3b: FALLBACK SE NON CI SONO I MODULI
# ─────────────────────────────────────────────
def process_pptx_fallback(pptx_path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(pptx_path))
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            title = ""
            if slide.shapes.title:
                title = slide.shapes.title.text.strip()
            lines.append(f"\n--- SLIDE {i}: {title} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t and t != title:
                            lines.append(t)
        print(f"    ✓ pptx fallback: {len(prs.slides)} slide")
        return "\n".join(lines)
    except Exception as e:
        print(f"    [ERRORE] pptx: {e}")
        return ""


# ─────────────────────────────────────────────
# STEP 4: PDF — estrazione a chunk di pagine
# ─────────────────────────────────────────────
def extract_pdf_pages(pdf_path: Path) -> list[dict]:
    """
    Estrae il testo del PDF pagina per pagina.
    Ritorna lista di dict: [{page: N, text: "..."}]
    """
    pages = []
    n_total_pages = 0
    try:
        import pdfplumber
        with pdfplumber.open(str(pdf_path)) as pdf:
            n_total_pages = len(pdf.pages)
            for i, page in enumerate(pdf.pages, 1):
                t = page.extract_text()
                if t and t.strip():
                    pages.append({"page": i, "text": t.strip()})
        print(f"    ✓ pdf pdfplumber: {len(pages)}/{n_total_pages} pagine con testo")
        if n_total_pages > 0 and not pages:
            print(f"    [WARN] {pdf_path.name}: nessun testo estratto — potrebbe essere un PDF scansionato (solo immagini).")
            print(f"           Per OCR su PDF scansionati: pip install pytesseract pdf2image")
        return pages
    except ImportError:
        pass
    try:
        import PyPDF2
        with open(pdf_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            n_total_pages = len(reader.pages)
            for i, page in enumerate(reader.pages, 1):
                t = page.extract_text()
                if t and t.strip():
                    pages.append({"page": i, "text": t.strip()})
        print(f"    ✓ pdf PyPDF2: {len(pages)}/{n_total_pages} pagine con testo")
        if n_total_pages > 0 and not pages:
            print(f"    [WARN] {pdf_path.name}: nessun testo estratto — potrebbe essere un PDF scansionato (solo immagini).")
            print(f"           Per OCR su PDF scansionati: pip install pytesseract pdf2image")
        return pages
    except ImportError:
        print("    [MANCANTE] pdfplumber — pip install pdfplumber")
        return []

# ──────────────────────────────
# PDF → singole pagine salvate
# ──────────────────────────────
def _ocr_pages_with_tesseract(images_dir: Path, page_images: dict) -> list[dict]:
    """
    OCR fallback per PDF scansionati (pdfplumber ha estratto 0 testo).
    Usa pytesseract sulle pagine già renderizzate come PNG da pdf_renderer.
    Ritorna lista [{page: N, text: "..."}] oppure [] se tesseract non installato.
    """
    try:
        import pytesseract
        from PIL import Image as _PILImage
    except ImportError:
        print("    [OCR PDF] pytesseract non installato — skip")
        print("              sudo apt install tesseract-ocr tesseract-ocr-ita")
        print("              pip install pytesseract")
        return []

    # Determina la lingua: mappa codici Whisper/BCP-47 → Tesseract
    _whisper_lang = os.environ.get("WHISPER_LANG", "").strip().lower()
    if _whisper_lang not in {"en", "it"}:
        _whisper_lang = ""  # ignora lingue non supportate
    _lang_map = {"it": "ita", "en": "eng"}
    tess_lang = _lang_map.get(_whisper_lang, "eng+ita")  # default: eng+ita (en prima)

    print(f"    [OCR PDF] pytesseract fallback — {len(page_images)} pagine, lang={tess_lang}")
    pages = []
    for page_num, img_filename in sorted(page_images.items()):
        img_path = images_dir / img_filename
        if not img_path.exists():
            continue
        try:
            text = pytesseract.image_to_string(_PILImage.open(img_path), lang=tess_lang)
            text = text.strip()
            if text:
                pages.append({"page": page_num, "text": text})
                print(f"    ✓ pag {page_num}: {len(text)} char")
        except Exception as e:
            print(f"    [OCR PDF] pagina {page_num}: {e}")

    print(f"    ✓ OCR completato: {len(pages)}/{len(page_images)} pagine con testo")
    return pages


def save_pdf_pages_as_txt(pdf_path: Path, output_dir: Path) -> list[Path]:
    """
    Estrae ogni pagina del PDF e la salva come file .txt nella cartella output_dir.
    Ritorna la lista dei file salvati.
    """
    output_dir.mkdir(parents=True, exist_ok=True)
    saved_files = []

    pages = extract_pdf_pages(pdf_path)
    for page in pages:
        page_num = page["page"]
        text     = page["text"]
        file_path = output_dir / f"slide_{page_num:03d}.txt"
        file_path.write_text(text, encoding="utf-8")
        saved_files.append(file_path)

    print(f"    ✓ Salvate {len(saved_files)} pagine in {output_dir}")
    return saved_files

def chunk_pdf_pages(pages: list[dict], chunk_size: int = 10) -> list[list[dict]]:
    """Divide le pagine in gruppi da chunk_size."""
    return [pages[i:i+chunk_size] for i in range(0, len(pages), chunk_size)]


def extract_pdf(pdf_path: Path) -> str:
    """Versione semplice per PDF piccoli o uso come testo extra."""
    pages = extract_pdf_pages(pdf_path)
    lines = [f"\n--- PAG {p['page']} ---\n{p['text']}" for p in pages]
    return "\n".join(lines)

def process_pdf_chunked(pdf_path: Path, output_dir: Path,
                         base_lesson_number: int, title: str,
                         skip_ai: bool = False,
                         chunk_size: int = 10,
                         subject_hint: str = None,
                         course_context_path: str = None) -> list[Path]:
    """
    Processa un PDF grande dividendolo in chunk di pagine.
    Ogni chunk diventa un lezione_NN.tex separato.
    Ritorna la lista dei .tex generati.
    """
    print(f"\n  [PDF grande] {pdf_path.name} — chunking ogni {chunk_size} pagine")
    images_dir = output_dir / "images"
    images_dir.mkdir(exist_ok=True)
    pages = extract_pdf_pages(pdf_path)
    if not pages:
        print("  [ERRORE] Nessun testo estratto dal PDF")
        return []

    total_pages = pages[-1]["page"] if pages else 0
    chunks      = chunk_pdf_pages(pages, chunk_size)
    print(f"  {total_pages} pagine totali → {len(chunks)} chunk da ~{chunk_size} pag")

    tex_files = []
    for idx, chunk in enumerate(chunks):
        if not chunk:
            continue
        lesson_num  = base_lesson_number + idx
        p_start     = chunk[0]["page"]
        p_end       = chunk[-1]["page"]
        chunk_title = f"{title} — pag. {p_start}–{p_end}"
        chunk_text  = "\n\n".join(
            f"[PAG {p['page']}]\n{p['text']}" for p in chunk
        )
        out_tex = output_dir / f"lezione_{lesson_num:02d}.tex"
        print(f"\n  Chunk {idx+1}/{len(chunks)}: pag {p_start}–{p_end}")

        if skip_ai:
            page_images, latex_skeleton = render_pdf_pages(
                pdf_path   = pdf_path,
                images_dir = images_dir,
                pages_data = chunk,
            ) if COLLEAGUE_MODULES else ({}, None)
            if latex_skeleton:
                content = (
                    f"\\section{{Lezione {lesson_num}: {_escape_latex(chunk_title)}}}\n"
                    f"\\label{{sec:lezione{lesson_num:02d}}}\n\n"
                    + latex_skeleton
                )
            else:
                content_lines = [
                    f"\\section{{Lezione {lesson_num}: {_escape_latex(chunk_title)}}}",
                    f"\\label{{sec:lezione{lesson_num:02d}}}\n",
                ]
                for p in chunk:
                    content_lines.append(f"\\subsection*{{Pagina {p['page']}}}")
                    for line in p["text"].split("\n"):
                        line = line.strip()
                        if line:
                            content_lines.append(_escape_latex(line) + "\n")
                content = "\n".join(content_lines)
        else:
            # ── Costruisce sources con la nuova struttura ──
            page_images, latex_skeleton = render_pdf_pages(
                pdf_path   = pdf_path,
                images_dir = images_dir,
                pages_data = chunk,
            ) if COLLEAGUE_MODULES else ({}, None)
            chunk_sources = {
                "has_audio": False,
                "scheletro": [{
                    "filename":    pdf_path.name,
                    "text":        chunk_text,
                    "pages":       len(chunk),
                    "latex":       latex_skeleton,
                    "page_images": page_images,
                }],
                "carne":    [],
                "supporto": [],
                "contorno": [],
            }
            content = generate_with_claude(
                lesson_number        = lesson_num,
                title                = chunk_title,
                sources              = chunk_sources,
                subject_hint         = subject_hint,
                course_context_path  = course_context_path,
                _progress_output_dir = output_dir,
            )
            if not content:
                content = (
                    f"\\section{{{_escape_latex(chunk_title)}}}\n"
                    f"\\label{{sec:lezione{lesson_num:02d}}}\n\n"
                )
                for p in chunk:
                    content += f"\\subsection*{{Pagina {p['page']}}}\n"
                    for line in p["text"].split("\n"):
                        line = line.strip()
                        if line:
                            content += _escape_latex(line) + "\n"

        write_lesson_tex(lesson_num, chunk_title, content,
                         [f"{pdf_path.name} pag.{p_start}-{p_end}"], out_tex)
        tex_files.append(out_tex)

    return tex_files
# ─────────────────────────────────────────────
# STEP 5: DOCX
# ─────────────────────────────────────────────
def extract_docx(docx_path: Path) -> str:
    """
    Estrae testo da DOCX preservando:
      - Paragrafi in ordine, con heading marcati come ## heading
      - Tabelle inline con delimitatore " | "
    Le tabelle sono inserite dopo i paragrafi per semplicità, ma vengono
    comunque integralmente estratte (evita perdita di informazioni strutturate).
    """
    try:
        from docx import Document
        from docx.oxml.ns import qn
        doc = Document(str(docx_path))
        parts = []

        # Percorri il body XML in ordine — paragrafi e tabelle interleaved
        for child in doc.element.body:
            local = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if local == "p":
                from docx.text.paragraph import Paragraph
                para = Paragraph(child, doc)
                txt = para.text.strip()
                if not txt:
                    continue
                if para.style and para.style.name.startswith("Heading"):
                    lvl = "".join(filter(str.isdigit, para.style.name)) or "1"
                    parts.append("#" * int(lvl) + " " + txt)
                else:
                    parts.append(txt)
            elif local == "tbl":
                from docx.table import Table as DocxTable
                tbl = DocxTable(child, doc)
                rows = []
                prev_cells: list = []          # evita righe duplicate (celle unite)
                for row in tbl.rows:
                    cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                    if cells == prev_cells:    # riga fantasma da merge verticale
                        continue
                    prev_cells = cells
                    if any(cells):
                        rows.append(" | ".join(cells))
                if rows:
                    parts.append("[TABELLA]\n" + "\n".join(rows))

        n_para  = sum(1 for p in doc.paragraphs if p.text.strip())
        n_table = len(doc.tables)
        print(f"    ✓ docx: {n_para} paragrafi, {n_table} tabelle")
        return "\n".join(parts)
    except ImportError:
        print("    [MANCANTE] python-docx — pip install python-docx")
        return ""
    except Exception as e:
        print(f"    [ERRORE] docx: {e}")
        return ""


# ─────────────────────────────────────────────
# STEP 6: GENERAZIONE LaTeX CON CLAUDE
# (con preprocessor integrato)
# ─────────────────────────────────────────────
try:
    from preprocessor import preprocess, NormalizedDocument, update_course_context, aligned_to_prompt
    PREPROCESSOR = True
except ImportError:
    PREPROCESSOR = False


def _trunc(text: str, max_chars: int, label: str = "") -> str:
    """
    Tronca il testo a max_chars per rispettare il budget token.
    Tronca al newline più vicino e aggiunge nota di troncamento.
    """
    if len(text) <= max_chars:
        return text
    cut = text[:max_chars]
    nl  = cut.rfind("\n")
    if nl > max_chars * 0.85:
        cut = cut[:nl]
    removed = len(text) - len(cut)
    tag = f" [{label}]" if label else ""
    print(f"    [WARN] Testo troncato{tag}: {len(text):,} → {len(cut):,} caratteri ({removed:,} rimossi)")
    return cut + f"\n\n[... {removed:,} caratteri omessi per rispettare il budget token ...]"


# Limiti per sezione (caratteri) — bilanciati per il contesto di claude-sonnet
_MAX_CARNE_CHARS    = 200_000   # trascrizione: la fonte più preziosa, taglio generoso
_MAX_SCHELETRO_CHARS = 160_000  # scheletro LaTeX: già strutturato, raramente supera
_MAX_SUPPORTO_CHARS  = 80_000   # PDF/DOCX di supporto: informazione supplementare
_MAX_CONTORNO_CHARS  = 40_000   # note informali: meno priorità


MAX_IMG_PX = 1920  # lato lungo massimo per le immagini PNG nello ZIP


def _resize_images_dir(images_dir: Path, max_px: int = MAX_IMG_PX) -> None:
    """Ridimensiona in-place tutti i PNG in images_dir se un lato supera max_px."""
    try:
        from PIL import Image
    except ImportError:
        return  # Pillow non disponibile — skip silenzioso
    resized = 0
    for png in images_dir.rglob("*.png"):
        try:
            with Image.open(png) as img:
                w, h = img.size
                if max(w, h) <= max_px:
                    continue
                scale = max_px / max(w, h)
                new_size = (int(w * scale), int(h * scale))
                img_resized = img.resize(new_size, Image.LANCZOS)
                img_resized.save(png, "PNG", optimize=True)
                resized += 1
        except Exception:
            pass
    if resized:
        print(f"    ✓ {resized} immagini ridimensionate a max {max_px}px")


def _clean_claude_output(text: str) -> str:
    """
    Rimuove artefatti comuni nelle risposte Claude prima che l'output
    venga scritto nel file .tex:
    - Blocchi markdown ```latex ... ``` o ``` ... ```
    - Testo introduttivo prima del primo comando LaTeX (\section, \chapter, ecc.)
    - Testo conclusivo dopo l'ultimo comando LaTeX
    """
    import re as _re

    # 1. Rimuovi code fences markdown (```latex ... ``` oppure ``` ... ```)
    text = _re.sub(r'^```(?:latex)?\s*\n', '', text, flags=_re.MULTILINE)
    text = _re.sub(r'\n```\s*$', '', text, flags=_re.MULTILINE)
    text = text.strip()

    # 2. Rimuovi testo prima del primo comando LaTeX rilevante
    #    (Claude a volte prepone "Ecco il LaTeX:", "Here is the code:", ecc.)
    _first = _re.search(r'\\(?:section|chapter|subsection|begin)\s*[\[{]', text)
    if _first and _first.start() > 0:
        prefix = text[:_first.start()].strip()
        if prefix:  # c'era davvero del testo spurio prima
            print(f"  [CLEANUP] Rimosso prefisso Claude ({len(prefix)} char): "
                  f"{prefix[:80]!r}{'…' if len(prefix)>80 else ''}")
        text = text[_first.start():]

    # 3. Rimuovi testo dopo l'ultimo } o \end{...} (note finali occasionali)
    _last = _re.search(r'(\\end\{[^}]+\}|^\})\s*$', text, flags=_re.MULTILINE)
    if _last:
        suffix = text[_last.end():].strip()
        if suffix:
            print(f"  [CLEANUP] Rimosso suffisso Claude ({len(suffix)} char): "
                  f"{suffix[:80]!r}{'…' if len(suffix)>80 else ''}")
            text = text[:_last.end()]

    return text.strip()


def generate_with_claude(lesson_number: int, title: str,
                          sources: dict,
                          subject_hint: str = None,
                          course_context_path: str = None,
                          _progress_output_dir: Path = None,
                          skip_ai: bool = False) -> str | None:
    """
    Genera LaTeX da Claude con prompt strutturato e gerarchia semantica.

    sources = {
        "has_audio":  bool,
        "scheletro":  [{"filename", "text", "latex", "slide_count", "slide_images"}],
        "carne":      [{"filename", "text"}],
        "supporto":   [{"filename", "text", "pages"}],
        "contorno":   [{"filename", "text"}],
    }
    """
    import urllib.request, urllib.error
    import json, os, time
    from pathlib import Path

    api_key = os.environ.get("ANTHROPIC_API_KEY")

    # ─────────────────────────────────────────
    # SYSTEM PROMPT — invariante
    # ─────────────────────────────────────────
    system_prompt = """Sei un esperto di LaTeX accademico. Trasforma le fonti fornite in un capitolo LaTeX professionale.

REGOLE OBBLIGATORIE:
1. Rispondi SOLO con codice LaTeX valido, iniziando da \\section{...}
2. NON includere \\documentclass, \\begin{document} o \\end{document}
3. Struttura: \\section{} > \\subsection{} > \\subsubsection{}
4. Aggiungi \\label{sec:...} a ogni section e subsection
5. Formule inline: $...$ — formule a display: \\begin{equation}...\\end{equation}
6. Liste: \\begin{itemize} o \\begin{enumerate}
7. Definizioni: \\begin{definition}...\\end{definition}
8. Teoremi: \\begin{theorem}...\\end{theorem}
9. Esempi: \\begin{example}...\\end{example}
10. Osservazioni e intuizioni del professore: \\begin{remark}...\\end{remark}
11. ZERO perdita di informazione concettuale dalla trascrizione: ogni spiegazione, esempio, intuizione e osservazione del professore deve apparire negli appunti — elimina solo le ripetizioni identiche e i riempitivi verbali ("allora", "quindi", "diciamo", ecc.)
12. Le spiegazioni orali che vanno OLTRE il contenuto delle slide sono il materiale più prezioso: preservale integralmente con lo stile del professore
13. Mantieni la terminologia tecnica originale del professore
14. Lo SCHELETRO contiene già formule in \\begin{equation}...\\end{equation} — mantienile senza modificarle
15. I blocchi \\begin{figure}...\\end{figure} nello SCHELETRO vanno mantenuti nella posizione esatta
16. Le tabelle del professore o del documento → \\begin{tabular} con header appropriato
17. Usa la lingua della trascrizione/slide per i contenuti (non forzare l'italiano se le fonti sono in un'altra lingua)
18. LIMITE ORALE — REGOLA ASSOLUTA: la lezione termina dove termina la spiegazione verbale del professore nella trascrizione. Se le slide contengono argomenti non ancora spiegati oralmente, NON includerli. Scrivere contenuto non detto dal professore è un errore grave
19. RACCORDO INTER-LEZIONE: se il CONTESTO CORSO indica "ultimo argomento trattato verbalmente", inizia da quel punto con un raccordo fluido di 1-2 righe. Non ripetere la spiegazione già svolta nella lezione precedente"""

    # ─────────────────────────────────────────
    # PREPROCESSOR — pulizia e contesto corso
    # ─────────────────────────────────────────
    course_context = ""
    subject        = subject_hint or "generico"

    if PREPROCESSOR:
        # Usa preprocessor solo per pulizia e contesto — non per assemblare il prompt
        carne_text      = "\n".join(s["text"] for s in sources["carne"])
        scheletro_raw   = "\n".join(s["text"] for s in sources["scheletro"])
        total_duration  = sum(s.get("duration_sec", 0) for s in sources["carne"]) or None

        doc = preprocess(
            transcript          = carne_text,
            slide_text          = scheletro_raw,
            extra_text          = "",
            title               = f"Lezione {lesson_number}: {title}",
            subject_hint        = subject_hint,
            course_context_path = course_context_path,
            lesson_number       = lesson_number,
            total_duration_sec  = total_duration,
        )
        subject        = doc.subject
        course_context = doc.context_prompt
        subject_instr  = doc.subject_prompt
        aligned        = doc.aligned_sections
    else:
        subject_instr = ""
        aligned       = []

    # ─────────────────────────────────────────
    # USER PROMPT — assemblaggio con gerarchia
    # ─────────────────────────────────────────
    sep  = "═" * 56
    sep2 = "─" * 56
    parts = []

    # ── Intestazione lezione ──
    has_scheletro = bool(sources["scheletro"])
    has_carne     = bool(sources["carne"])
    has_supporto  = bool(sources["supporto"])
    has_contorno  = bool(sources["contorno"])

    fonti_str = []
    if has_scheletro: fonti_str.append("scheletro")
    if has_carne:     fonti_str.append("carne")
    if has_supporto:  fonti_str.append("supporto")
    if has_contorno:  fonti_str.append("contorno")

    parts.append(
        f"{sep}\n"
        f"  LEZIONE {lesson_number}: {title}\n"
        f"  Materia: {subject} | Fonti: {', '.join(fonti_str)}\n"
        f"{sep}"
    )

    # ── Contesto corso ──
    if course_context:
        parts.append(f"{sep}\n  CONTESTO CORSO\n{sep}\n{course_context}")

    # ── Istruzioni materia ──
    if subject_instr:
        parts.append(f"{sep}\n  ISTRUZIONI MATERIA\n{sep}\n{subject_instr}")

    # ── SCHELETRO ──
    if has_scheletro:
        parts.append(f"{sep}\n  FONTE: SCHELETRO\n{sep}")

        for entry in sources["scheletro"]:
            filename    = entry["filename"]
            latex       = entry.get("latex")
            text        = entry.get("text", "")
            slide_count = entry.get("slide_count", "?")
            pages       = entry.get("pages", "?")

            is_pptx = filename.lower().endswith(".pptx")
            is_pdf  = filename.lower().endswith(".pdf")
            is_docx = filename.lower().endswith(".docx")

            if is_pptx:
                meta = f"File: {filename} | {slide_count} slide"
                role = (
                    "Struttura UFFICIALE della lezione generata dalle slide.\n"
                    "REGOLE:\n"
                    "  • Ogni \\subsection corrisponde a una slide — non cambiare questa struttura\n"
                    "  • I blocchi \\begin{figure} contengono le immagini delle slide — mantienili nella posizione esatta\n"
                    "  • Le formule in \\begin{equation} sono già in LaTeX verificato — mantienile senza modifiche\n"
                    "  • Arricchisci il contenuto testuale con la CARNE (trascrizione)"
                )
                raw     = latex if latex else text
                content = _trunc(raw, _MAX_SCHELETRO_CHARS, filename)

            elif is_pdf:
                meta = f"File: {filename} | {pages} pagine"
                role = (
                    "Documento PDF che funge da SCHELETRO perché è presente la trascrizione audio.\n"
                    "REGOLE:\n"
                    "  • Usa la struttura del documento (titoli, sezioni) come guida per \\subsection\n"
                    "  • Ogni [PAG N] è una pagina del documento — usala come unità strutturale\n"
                    "  • Arricchisci con la CARNE (trascrizione audio)"
                )
                raw     = entry.get("latex") or text
                content = _trunc(raw, _MAX_SCHELETRO_CHARS, filename)

            elif is_docx:
                meta = f"File: {filename}"
                role = (
                    "Documento Word che funge da SCHELETRO perché è presente la trascrizione audio.\n"
                    "REGOLE:\n"
                    "  • Usa i paragrafi e le tabelle come guida per la struttura \\subsection\n"
                    "  • Le tabelle [TABELLA] vanno convertite in \\begin{tabular}\n"
                    "  • Arricchisci con la CARNE (trascrizione audio)"
                )
                content = _trunc(text, _MAX_SCHELETRO_CHARS, filename)

            else:
                meta    = f"File: {filename}"
                role    = "Documento strutturale della lezione."
                content = _trunc(text, _MAX_SCHELETRO_CHARS, filename)

            parts.append(
                f"{sep2}\n{meta}\n\n{role}\n{sep2}\n{content}"
            )

    # ── CARNE ──
    if has_carne:
        parts.append(f"{sep}\n  FONTE: CARNE (voce del professore)\n{sep}")

        # Se abbiamo l'allineamento slide↔trascrizione, usalo — è molto più utile per Claude
        if aligned and has_scheletro:
            aligned_text = aligned_to_prompt(aligned)
            parts.append(
                f"{sep2}\n"
                f"Trascrizione ALLINEATA per slide — voce del professore mappata sulla struttura\n"
                f"REGOLE:\n"
                f"  • Ogni blocco [SPIEGAZIONE ORALE] corrisponde alla slide soprastante\n"
                f"  • Le ripetizioni di un concetto indicano importanza — enfatizzalo\n"
                f"  • Gli esempi verbali non nello scheletro → \\begin{{example}}\n"
                f"  • Le frasi 'quindi', 'in altre parole', 'ricordate' → spiegazioni chiave\n"
                f"{sep2}\n{aligned_text}"
            )
        else:
            # REGOLE scritte una volta nella header — non ripetute per ogni file
            parts.append(
                f"Trascrizione della VOCE DEL PROFESSORE durante la lezione.\n"
                f"REGOLE:\n"
                f"  • Integra le spiegazioni nelle \\subsection corrispondenti dello SCHELETRO\n"
                f"  • Le ripetizioni di un concetto indicano importanza — enfatizzalo\n"
                f"  • Gli esempi verbali non presenti nello scheletro → \\begin{{example}}\n"
                f"  • Le frasi 'quindi', 'in altre parole', 'ricordate' → spiegazioni chiave\n"
                f"  • I timestamp [MM:SS] indicano la progressione temporale"
            )
            for entry in sources["carne"]:
                filename = entry["filename"]
                text     = _trunc(entry["text"], _MAX_CARNE_CHARS, filename)
                parts.append(f"{sep2}\nFile: {filename}\n{sep2}\n{text}")

    # ── SUPPORTO ──
    if has_supporto:
        parts.append(f"{sep}\n  FONTE: SUPPORTO (materiale di riferimento)\n{sep}")

        for entry in sources["supporto"]:
            filename = entry["filename"]
            text     = _trunc(entry["text"], _MAX_SUPPORTO_CHARS, filename)
            pages    = entry.get("pages", "")
            meta     = f"File: {filename}" + (f" | {pages} pagine" if pages else "")
            parts.append(
                f"{sep2}\n"
                f"{meta}\n\n"
                f"Materiale di SUPPORTO — non è la struttura della lezione.\n"
                f"REGOLE:\n"
                f"  • Usalo per arricchire definizioni formali dove scheletro/carne sono sintetici\n"
                f"  • Le tabelle [TABELLA] → \\begin{{tabular}} con intestazioni chiare\n"
                f"  • Non cambiare la struttura \\subsection per adattarla a questo documento\n"
                f"{sep2}\n{text}"
            )

    # ── CONTORNO ──
    if has_contorno:
        parts.append(f"{sep}\n  FONTE: CONTORNO (note informali)\n{sep}")

        for entry in sources["contorno"]:
            filename = entry["filename"]
            text     = _trunc(entry["text"], _MAX_CONTORNO_CHARS, filename)
            parts.append(
                f"{sep2}\n"
                f"File: {filename}\n\n"
                f"Note informali — peso MINORE rispetto alle altre fonti.\n"
                f"Usa solo per dettagli non coperti altrove.\n"
                f"{sep2}\n{text}"
            )

    # ── ISTRUZIONI DI SINTESI — adattive ──
    parts.append(f"{sep}\n  ISTRUZIONI DI SINTESI\n{sep}")
    instructions = _build_synthesis_instructions(
        lesson_number, title,
        has_scheletro, has_carne, has_supporto,
        sources,
    )
    parts.append(instructions)

    user_prompt = "\n\n".join(parts)

    # ─────────────────────────────────────────
    # IMMAGINI — raccolta path + base64
    # (prima del debug così il salvataggio include tutto)
    # ─────────────────────────────────────────
    import base64, os as _os
    _MAX_IMAGES = 20   # limite API Claude per chiamata

    images_dir   = (_progress_output_dir / "images") if _progress_output_dir else None
    image_paths: list[Path] = []   # per debug
    image_blocks: list[dict] = []  # per API

    if images_dir and images_dir.exists():
        # 1. Slide PPTX — priorità massima (contenuto visivo irriducibile a testo)
        for entry in sources["scheletro"]:
            if not entry["filename"].lower().endswith(".pptx"):
                continue
            for slide_num in sorted((entry.get("slide_images") or {}).keys()):
                if len(image_blocks) >= _MAX_IMAGES:
                    break
                img_path = images_dir / entry["slide_images"][slide_num]
                if img_path.exists():
                    data = base64.b64encode(img_path.read_bytes()).decode()
                    image_paths.append(img_path)
                    image_blocks.append({
                        "type":   "image",
                        "source": {"type": "base64", "media_type": "image/png", "data": data},
                    })

        # 2. Pagine PDF — solo se non c'è audio e slot liberi
        if not sources["carne"]:
            for entry in sources["scheletro"]:
                if not entry["filename"].lower().endswith(".pdf"):
                    continue
                for page_num in sorted((entry.get("page_images") or {}).keys()):
                    if len(image_blocks) >= _MAX_IMAGES:
                        break
                    img_path = images_dir / entry["page_images"][page_num]
                    if img_path.exists():
                        data = base64.b64encode(img_path.read_bytes()).decode()
                        image_paths.append(img_path)
                        image_blocks.append({
                            "type":   "image",
                            "source": {"type": "base64", "media_type": "image/png", "data": data},
                        })

    # ─────────────────────────────────────────
    # DEBUG — salva TUTTO ciò che andrebbe a Claude
    # ─────────────────────────────────────────
    debug_dir = (_progress_output_dir / "debug") if _progress_output_dir else Path("debug")
    debug_dir.mkdir(parents=True, exist_ok=True)

    # 1. Testo: system + user prompt
    debug_path = debug_dir / f"prompt_lezione_{lesson_number:02d}.txt"
    img_section = ""
    if image_paths:
        img_lines = "\n".join(f"  [{i+1:02d}] {p.name}" for i, p in enumerate(image_paths))
        img_section = f"\n\n=== IMMAGINI ALLEGATE ({len(image_paths)}) ===\n{img_lines}"
    debug_path.write_text(
        f"=== SYSTEM ===\n{system_prompt}\n\n=== USER ===\n{user_prompt}{img_section}",
        encoding="utf-8"
    )

    # 2. Immagini: symlink ordinati in debug/images_lezione_NN/
    if image_paths:
        dbg_img_dir = debug_dir / f"images_lezione_{lesson_number:02d}"
        dbg_img_dir.mkdir(exist_ok=True)
        # Rimuovi symlink precedenti (riesecuzione)
        for old in dbg_img_dir.iterdir():
            if old.is_symlink():
                old.unlink()
        for i, src in enumerate(image_paths, 1):
            link = dbg_img_dir / f"{i:02d}_{src.name}"
            try:
                link.symlink_to(src.resolve())
            except Exception:
                pass  # filesystem senza symlink: ignora

    est_tokens = (len(system_prompt) + len(user_prompt)) // 4
    print(f"\n  [DEBUG] Prompt → {debug_path.resolve()}")
    if image_paths:
        print(f"  [DEBUG] Immagini → {debug_dir / f'images_lezione_{lesson_number:02d}'}/ ({len(image_paths)} file)")
    print(f"  [DEBUG] ~{est_tokens:,} token stimati (testo) + {len(image_paths)} immagini")
    print(f"  [DEBUG] Fonti: scheletro={len(sources['scheletro'])} "
          f"carne={len(sources['carne'])} "
          f"supporto={len(sources['supporto'])} "
          f"contorno={len(sources['contorno'])}")

    if skip_ai:
        print("  [SKIP] --skip-ai: prompt salvato, chiamata API saltata")
        return None

    if not api_key:
        print("  [SKIP] ANTHROPIC_API_KEY non impostata")
        return None

    # ─────────────────────────────────────────
    # Assembla content multimodale
    # ─────────────────────────────────────────
    if image_blocks:
        print(f"  [VISION] {len(image_blocks)} immagini allegate alla chiamata API")
        user_content = [
            {"type": "text", "text": f"Hai a disposizione {len(image_blocks)} immagini delle slide/pagine della lezione, nell'ordine in cui appaiono nel prompt seguente.\n"},
            *image_blocks,
            {"type": "text", "text": user_prompt},
        ]
    else:
        user_content = user_prompt

    # ─────────────────────────────────────────
    # CHIAMATA API
    # ─────────────────────────────────────────
    # system come lista di blocchi con cache_control — il testo è identico per ogni lezione
    # → Anthropic lo cachea dopo la prima chiamata (~70% risparmio su quei token)
    payload = json.dumps({
        "model":      CONFIG["claude_model"],
        "max_tokens": CONFIG["claude_max_tokens"],
        "system": [
            {
                "type":          "text",
                "text":          system_prompt,
                "cache_control": {"type": "ephemeral"},
            }
        ],
        "messages": [{"role": "user", "content": user_content}],
    }).encode("utf-8")

    req = urllib.request.Request(
        "https://api.anthropic.com/v1/messages",
        data    = payload,
        headers = {
            "Content-Type":      "application/json",
            "x-api-key":         api_key,
            "anthropic-version": "2023-06-01",
            "anthropic-beta":    "prompt-caching-2024-07-31",
        },
        method = "POST",
    )

    print(f"  Claude API: lezione {lesson_number} ...")
    if _progress_output_dir:
        _report_progress(_progress_output_dir, 70,
                         "Claude — Generazione LaTeX",
                         f"Lezione {lesson_number}: {title}")

    # Retry con backoff esponenziale: 3 tentativi su errori transitori
    _RETRY_CODES  = {429, 529}   # rate limit / overloaded
    _MAX_ATTEMPTS = 3
    _BACKOFF_BASE = 20           # secondi: 20s, 40s

    for attempt in range(1, _MAX_ATTEMPTS + 1):
        t0 = time.time()
        try:
            with urllib.request.urlopen(req, timeout=300) as resp:
                data  = json.loads(resp.read())
                content_blocks = data.get("content") or []
                if not content_blocks:
                    raise ValueError(f"Risposta Claude vuota o malformata: {str(data)[:200]}")
                latex = _clean_claude_output(content_blocks[0].get("text", ""))
                usage = data.get("usage", {})
                cache_hit  = usage.get("cache_read_input_tokens", 0)
                cache_miss = usage.get("cache_creation_input_tokens", 0)
                cache_info = (f", cache={'hit' if cache_hit else 'miss'} "
                              f"(r:{cache_hit} w:{cache_miss})") if (cache_hit or cache_miss) else ""
                print(f"  ✓ Claude: {time.time()-t0:.1f}s, {len(latex):,} chars{cache_info}")

                if PREPROCESSOR and course_context_path:
                    update_course_context(
                        context_path  = course_context_path,
                        lesson_number = lesson_number,
                        lesson_title  = title,
                        latex_content = latex,
                    )
                return latex

        except urllib.error.HTTPError as e:
            body = e.read().decode(errors="replace")[:300]
            if e.code in _RETRY_CODES and attempt < _MAX_ATTEMPTS:
                wait = _BACKOFF_BASE * attempt
                print(f"  [retry {attempt}/{_MAX_ATTEMPTS}] Claude HTTP {e.code} — riprovo tra {wait}s")
                time.sleep(wait)
                continue
            print(f"  [ERRORE] Claude HTTP {e.code}: {body}")
            return None
        except (TimeoutError, OSError) as e:
            if attempt < _MAX_ATTEMPTS:
                wait = _BACKOFF_BASE * attempt
                print(f"  [retry {attempt}/{_MAX_ATTEMPTS}] Claude timeout/rete ({e}) — riprovo tra {wait}s")
                time.sleep(wait)
                continue
            print(f"  [ERRORE] Claude: {e}")
            return None
        except Exception as e:
            print(f"  [ERRORE] Claude: {e}")
            return None

    return None


# ─────────────────────────────────────────────
# ISTRUZIONI ADATTIVE — cambiano in base alle fonti
# ─────────────────────────────────────────────

def _build_synthesis_instructions(lesson_number: int, title: str,
                                   has_scheletro: bool, has_carne: bool,
                                   has_supporto: bool, sources: dict) -> str:
    lines = [
        f"Genera: \\section{{Lezione {lesson_number}: {title}}} e tutto il contenuto.",
        "",
    ]

    # Caso 1: scheletro PPTX + carne audio — caso ideale
    has_pptx_skeleton = any(
        s["filename"].lower().endswith(".pptx")
        for s in sources.get("scheletro", [])
    )

    if has_pptx_skeleton and has_carne:
        lines += [
            "CASO: slide + audio (caso ideale)",
            "1. Parti dallo SCHELETRO — mantieni ogni \\subsection e ogni \\begin{figure} nella posizione esatta",
            "2. Per ogni \\subsection espandi con la SPIEGAZIONE COMPLETA del professore dalla CARNE:",
            "   - Trascrivi integralmente la spiegazione del professore per quella slide",
            "   - Preserva il filo logico e lo stile didattico originale",
            "   - Rimuovi solo i riempitivi verbali puri (\"allora\", \"quindi\", ripetizioni identiche)",
            "3. Se scheletro e carne si contraddicono → privilegia lo SCHELETRO per i fatti, la CARNE per le spiegazioni",
            "4. Esempi verbali del professore non nello scheletro → aggiungi come \\begin{example}",
            "5. Osservazioni, intuizioni, avvertenze del professore → aggiungi come \\begin{remark}",
            "6. Ogni subsection deve contenere la spiegazione orale completa, non solo i bullet point della slide",
            "7. Formule [FORMULA_OMML] → copia direttamente senza modifiche",
            "8. Formule pronunciate nella trascrizione → converti tu in LaTeX",
        ]

    # Caso 2: solo scheletro PPTX, no audio
    elif has_pptx_skeleton and not has_carne:
        lines += [
            "CASO: solo slide (nessun audio)",
            "1. Mantieni ogni \\subsection e ogni \\begin{figure} nella posizione esatta",
            "2. Espandi ogni punto della slide con elaborazione accademica approfondita",
            "3. Aggiungi contesto, motivazioni e implicazioni per ogni concetto",
            "4. Formule [FORMULA_OMML] → copia direttamente senza modifiche",
        ]

    # Caso 3: scheletro PDF/DOCX + carne audio
    elif has_scheletro and has_carne:
        lines += [
            "CASO: documento + audio",
            "1. Usa la struttura del documento come guida per \\subsection",
            "2. Per ogni sezione integra la SPIEGAZIONE COMPLETA del professore dalla trascrizione:",
            "   - Preserva ogni spiegazione, esempio e osservazione del professore",
            "   - Rimuovi solo i riempitivi verbali puri",
            "3. Se documento e trascrizione si contraddicono → privilegia il documento per i fatti",
            "4. Esempi verbali → \\begin{example}, osservazioni → \\begin{remark}",
        ]

    # Caso 4: solo audio, niente scheletro
    elif not has_scheletro and has_carne:
        lines += [
            "CASO: solo audio (nessuno scheletro)",
            "1. Struttura autonomamente identificando i macro-argomenti nella trascrizione",
            "2. Ogni cambio di argomento → nuova \\subsection",
            "3. Preserva integralmente le spiegazioni del professore — zero perdita di informazione",
            "4. Segui l'ordine cronologico della lezione",
            "5. Esempi verbali → \\begin{example}, osservazioni → \\begin{remark}",
            "6. Formule pronunciate → converti in LaTeX",
        ]

    # Caso 5: solo documenti, niente audio
    elif has_scheletro and not has_carne:
        lines += [
            "CASO: solo documenti (nessun audio)",
            "1. Riformatta il documento in LaTeX strutturato",
            "2. Ogni sezione/capitolo del documento → \\subsection",
            "3. Espandi dove necessario per chiarezza accademica",
        ]

    if has_supporto:
        lines += [
            "",
            "SUPPORTO disponibile:",
            "• Usa per arricchire definizioni formali dove scheletro/carne sono sintetici",
            "• Non cambiare la struttura \\subsection per adattarla al supporto",
        ]

    lines += [
        "",
        "QUALITÀ ATTESA:",
        "• Struttura gerarchica chiara e navigabile",
        "• Spiegazioni complete del professore integrate per ogni concetto, non solo i bullet point",
        "• Nessuna informazione concettuale della trascrizione va persa",
        "• Formule matematiche corrette e leggibili",
        "• Pronto per compilazione con pdflatex",
    ]

    return "\n".join(lines)# ─────────────────────────────────────────────
# STEP 7: LaTeX STRUTTURATO SENZA AI
# Usa i dati del collega quando disponibili
# ─────────────────────────────────────────────
def build_fallback_latex(lesson_number: int, title: str,
                          slides, transcript: str,
                          slide_text: str, extra_text: str,
                          slide_images: dict = None) -> str:
    """
    Costruisce LaTeX strutturato senza Claude.

    Ogni \subsection corrisponde a una slide e include:
      1. \begin{figure} con il PNG della slide (se disponibile)
      2. Contenuto testuale della slide (testo + formule OMML)

    slide_images: {slide_number: "slide_001.png", ...}
                  prodotto da render_slide_images()
    """
    slide_images = slide_images or {}

    if COLLEAGUE_MODULES and slides:
        esc = _escape_latex
        parts = []
        parts.append(f"\\section{{Lezione {lesson_number}: {esc(title)}}}")
        parts.append(f"\\label{{sec:lezione{lesson_number:02d}}}\n")

        for slide in slides:
            sec_title = slide.title.strip() if slide.title.strip() else f"Slide {slide.slide_number}"
            parts.append(f"\n\\subsection{{{esc(sec_title)}}}")
            parts.append(f"\\label{{subsec:slide{lesson_number:02d}-{slide.slide_number}}}\n")

            # ── Figura slide PNG (se disponibile) ──
            img_filename = slide_images.get(slide.slide_number)
            if img_filename:
                parts.append(slide_figure_latex(
                    img_filename = img_filename,
                    slide_number = slide.slide_number,
                    caption      = sec_title,
                ))

            # ── Contenuto testuale + formule ──
            for obj in slide.objects:
                if obj.obj_type == "text":
                    if obj.content.strip() == slide.title.strip():
                        continue
                    lines = obj.content.strip().split("\n")
                    has_bullets = any(
                        l.strip().startswith(("•", "-", "*", "–"))
                        for l in lines
                    )
                    if has_bullets:
                        parts.append("\\begin{itemize}")
                        for l in lines:
                            l = l.strip()
                            if l:
                                parts.append(f"  \\item {esc(l.lstrip('•-*– ').strip())}")
                        parts.append("\\end{itemize}\n")
                    else:
                        body = "\n".join(esc(l) for l in lines if l.strip())
                        if body:
                            parts.append(body + "\n")

                elif obj.obj_type == "omml_formula":
                    f = getattr(obj, "latex_result", "")
                    if f and not f.startswith("%"):
                        parts.append("\\begin{equation}")
                        parts.append(f)
                        parts.append("\\end{equation}\n")
                    elif f:
                        parts.append(f"% OMML parziale:\n% {f}\n")

                elif obj.obj_type == "image":
                    f = getattr(obj, "latex_result", "")
                    img = obj.content
                    if f and f.strip():
                        # Formula riconosciuta da pix2tex
                        parts.append("\\begin{equation}")
                        parts.append(f)
                        parts.append("\\end{equation}\n")
                    else:
                        # Immagine embedded normale
                        parts.append("\\begin{figure}[H]")
                        parts.append("  \\centering")
                        parts.append(f"  \\includegraphics[width=0.8\\textwidth]{{{img}}}")
                        parts.append("\\end{figure}\n")

        # Trascrizione in fondo se disponibile
        if transcript:
            parts.append("\n\\subsection{Trascrizione Audio}")
            parts.append("\\begin{quote}")
            for line in transcript.split("\n")[:60]:
                parts.append(_escape_latex(line))
            if len(transcript.split("\n")) > 60:
                parts.append("\\emph{[trascrizione troncata]}")
            parts.append("\\end{quote}\n")

        return "\n".join(parts)

    # ── Fallback puro (nessun modulo collega) ──
    def esc(t):
        for a, b in [("\\", "\\textbackslash{}"), ("&", "\\&"), ("%", "\\%"),
                     ("$", "\\$"), ("#", "\\#"), ("{", "\\{"), ("}", "\\}")]:
            t = t.replace(a, b)
        return t

    parts = []
    parts.append(f"\\section{{Lezione {lesson_number}: {esc(title)}}}\n")

    if slide_text:
        current_slide_num = None
        for line in slide_text.split("\n"):
            line = line.strip()
            if line.startswith("--- SLIDE"):
                # Estrai numero slide per recuperare il PNG
                import re
                m = re.match(r"--- SLIDE (\d+)", line)
                current_slide_num = int(m.group(1)) if m else None
                parts.append(f"\n\\subsection{{{esc(line)}}}\n")
                # Aggiungi figura se disponibile
                if current_slide_num and current_slide_num in slide_images:
                    parts.append(slide_figure_latex(
                        img_filename = slide_images[current_slide_num],
                        slide_number = current_slide_num,
                    ))
            elif line:
                parts.append(esc(line) + "\n")

    if transcript:
        parts.append("\\subsection{Trascrizione Audio}\n\\begin{quote}")
        for line in transcript.split("\n")[:50]:
            parts.append(esc(line))
        parts.append("\\end{quote}\n")

    if extra_text:
        parts.append("\\subsection{Note Aggiuntive}\n\\begin{quote}")
        for line in extra_text.split("\n")[:30]:
            parts.append(esc(line))
        parts.append("\\end{quote}\n")

    return "\n".join(parts)
# ─────────────────────────────────────────────
# SCRITTURA lezione_NN.tex
# ─────────────────────────────────────────────
def write_lesson_tex(lesson_number: int, title: str,
                     content: str, sources: list, out_path: Path):
    header = (
        f"% {'='*60}\n"
        f"% Lezione {lesson_number:02d}: {title}\n"
        f"% Generato: {datetime.now().strftime('%Y-%m-%d %H:%M')}\n"
        f"% Fonti: {', '.join(sources) if sources else '—'}\n"
        f"% {'='*60}\n\n"
    )
    out_path.write_text(header + content, encoding="utf-8")
    kb = out_path.stat().st_size // 1024 + 1
    print(f"    ✓ {out_path.name}  ({kb} KB)")


# ─────────────────────────────────────────────
# STATE — numerazione persistente per corso
# ─────────────────────────────────────────────
_STATE_FILE = "state.json"


def load_state(output_dir: Path) -> dict:
    """
    Legge state.json dalla cartella output del corso.
    Struttura:
      {
        "course_title": "...",
        "subject":      "...",
        "next_lesson":  N,
        "lessons": [
          {"number": 1, "source_dir": "...", "tex_file": "lezione_01.tex",
           "processed_at": "ISO8601"}
        ]
      }
    """
    path = output_dir / _STATE_FILE
    if path.exists():
        try:
            import json as _json
            return _json.loads(path.read_text(encoding="utf-8"))
        except Exception:
            pass
    return {"course_title": "", "subject": None, "next_lesson": 1, "lessons": []}


def save_state(output_dir: Path, state: dict):
    """Salva state.json nella cartella output del corso (scrittura atomica)."""
    import json as _json
    path = output_dir / _STATE_FILE
    tmp  = path.with_suffix(".json.tmp")
    tmp.write_text(_json.dumps(state, indent=2, ensure_ascii=False), encoding="utf-8")
    tmp.replace(path)   # atomico: nessun file corrotto se il processo viene killato


# ─────────────────────────────────────────────
# PROGRESS — comunicazione real-time col server
# ─────────────────────────────────────────────

def _report_progress(output_dir: Path, progress: int,
                     current_step: str, detail: str = ""):
    """
    Scrive progress.json nella output_dir.
    Il server FastAPI lo legge ad ogni GET /job/{job_id}.
    Silenzioso in caso di errore (non deve mai crashare la pipeline).
    """
    import json as _json
    try:
        data = {
            "progress":     max(0, min(100, progress)),
            "current_step": current_step,
            "detail":       detail,
        }
        (output_dir / "progress.json").write_text(
            _json.dumps(data), encoding="utf-8"
        )
    except Exception:
        pass


# ─────────────────────────────────────────────
# GENERAZIONE main.tex
# ─────────────────────────────────────────────
MAIN_TEMPLATE = r"""\documentclass[12pt,a4paper]{{report}}

% ---------------------------------------------------------
% ENCODING & LINGUA
% ---------------------------------------------------------
\usepackage[utf8]{{inputenc}}
\usepackage[T1]{{fontenc}}
\usepackage[english,italian]{{babel}}
% Lingua principale: impostata da \selectlanguage in titlepage
\usepackage{{lmodern}}

% ---------------------------------------------------------
% MATEMATICA
% ---------------------------------------------------------
\usepackage{{amsmath,amsfonts,amssymb,amsthm,mathtools}}
\usepackage{{gensymb}}

% ---------------------------------------------------------
% LAYOUT & GEOMETRIA
% ---------------------------------------------------------
\usepackage[margin=2.5cm]{{geometry}}
\usepackage{{microtype}}

% ---------------------------------------------------------
% IMMAGINI & FIGURE
% ---------------------------------------------------------
\usepackage{{graphicx,float,subcaption}}
\graphicspath{{{{{images_path}}}}}

% ---------------------------------------------------------
% COLORI & GRAFICA
% ---------------------------------------------------------
\usepackage{{xcolor}}
\usepackage{{tikz}}
\usetikzlibrary{{positioning,calc,arrows.meta,shapes,shapes.geometric}}
\tikzstyle{{block}} = [draw, rectangle, minimum height=1.2cm, minimum width=1.8cm]
\tikzstyle{{sum}}   = [draw, circle, inner sep=0pt, minimum size=6mm]
\tikzstyle{{input}} = [coordinate]

% ---------------------------------------------------------
% LISTE & SIMBOLI
% ---------------------------------------------------------
\usepackage{{enumitem}}
\usepackage{{pifont}}
\usepackage{{newunicodechar}}
\newunicodechar{{✓}}{{\ding{{51}}}}
\newunicodechar{{✗}}{{\ding{{55}}}}

% ---------------------------------------------------------
% CODICE
% ---------------------------------------------------------
\usepackage{{listings}}
\lstset{{
  basicstyle=\ttfamily\small,
  breaklines=true,
  frame=single,
  backgroundcolor=\color{{gray!10}},
  numbers=left,
  numberstyle=\tiny\color{{gray}},
  keywordstyle=\color{{blue!70!black}}\bfseries,
  commentstyle=\color{{green!50!black}}\itshape,
  stringstyle=\color{{orange!80!black}},
}}

% ---------------------------------------------------------
% AMBIENTI COLORATI (tcolorbox)
% ---------------------------------------------------------
\usepackage[most]{{tcolorbox}}

\tcbset{{
  theorembase/.style={{
    enhanced, breakable,
    fonttitle=\bfseries,
    separator sign={{~—~}},
  }}
}}

\newtcbtheorem[number within=chapter]{{theorem}}{{Teorema}}{{
  theorembase,
  colback=blue!5!white, colframe=blue!50!black,
}}{{thm}}

\newtcbtheorem[use counter from=theorem]{{definition}}{{Definizione}}{{
  theorembase,
  colback=green!5!white, colframe=green!45!black,
}}{{def}}

\newtcbtheorem[use counter from=theorem]{{example}}{{Esempio}}{{
  theorembase,
  colback=orange!5!white, colframe=orange!60!black,
}}{{ex}}

\newtcbtheorem[use counter from=theorem]{{remark}}{{Osservazione}}{{
  theorembase,
  colback=gray!8!white, colframe=gray!50!black,
}}{{rem}}

\newtcbtheorem[use counter from=theorem]{{lemma}}{{Lemma}}{{
  theorembase,
  colback=purple!5!white, colframe=purple!50!black,
}}{{lem}}

% ---------------------------------------------------------
% HEADER / FOOTER
% ---------------------------------------------------------
\usepackage{{fancyhdr}}
\pagestyle{{fancy}}
\fancyhf{{}}
\lhead{{\leftmark}}
\rhead{{{title}}}
\cfoot{{\thepage}}
\renewcommand{{\headrulewidth}}{{0.4pt}}
\renewcommand{{\footrulewidth}}{{0pt}}

% ---------------------------------------------------------
% LINK
% ---------------------------------------------------------
\usepackage{{hyperref}}
\hypersetup{{
  colorlinks = true,
  linkcolor  = blue!70!black,
  urlcolor   = blue!70!black,
  citecolor  = green!50!black,
}}

% ---------------------------------------------------------
% INIZIO DOCUMENTO
% ---------------------------------------------------------
\begin{{document}}

\begin{{titlepage}}
\centering
\vspace*{{3cm}}
{{\Huge\bfseries {title}\par}}
\vspace{{1cm}}
{{\large Appunti generati automaticamente\par}}
\vspace{{0.5cm}}
{{\large {date}\par}}
\vfill
{{\normalsize Generato con \texttt{{appunti\_ai}} + Claude (Anthropic)\par}}
\end{{titlepage}}

\tableofcontents
\clearpage

{includes}

\end{{document}}
"""

def _latex_escape_title(t: str) -> str:
    """Escape caratteri speciali LaTeX nel titolo (usato solo in generate_main_tex)."""
    if COLLEAGUE_MODULES:
        return _escape_latex(t)
    # Fallback inline se _escape_latex non è disponibile
    for a, b in [
        ("\\", "\\textbackslash{}"),
        ("&",  "\\&"),
        ("%",  "\\%"),
        ("$",  "\\$"),
        ("#",  "\\#"),
        ("_",  "\\_"),
        ("^",  "\\^{}"),
        ("{",  "\\{"),
        ("}",  "\\}"),
        ("~",  "\\~{}"),
    ]:
        t = t.replace(a, b)
    return t


def generate_main_tex(title: str, lesson_files: list, output_dir: Path,
                      lang: str = "italian") -> Path:
    title_tex = _latex_escape_title(title)
    includes = "\n".join(f"\\include{{{f.stem}}}" for f in sorted(lesson_files))
    # Seleziona lingua principale in base alla lingua rilevata del corso
    _babel_lang = {
        "it": "italian", "en": "english", "fr": "french",
        "de": "ngerman", "es": "spanish", "pt": "portuguese",
    }.get(lang[:2] if lang else "it", "italian")
    content = MAIN_TEMPLATE.format(
        title=title_tex,
        date=datetime.now().strftime("%B %Y"),
        images_path=CONFIG["images_subdir"] + "/",
        includes=includes,
    ).replace(
        "% Lingua principale: impostata da \\selectlanguage in titlepage",
        f"% Lingua principale: impostata da \\selectlanguage in titlepage\n\\selectlanguage{{{_babel_lang}}}",
    )
    main_path = output_dir / "main.tex"
    main_path.write_text(content, encoding="utf-8")
    print(f"  ✓ main.tex  ({len(lesson_files)} lezioni incluse)")
    return main_path


# ─────────────────────────────────────────────
# HELPER: raccoglitori per i 6 tipi di sorgente
# ─────────────────────────────────────────────

import re as _re_pl

_TS_PAT  = _re_pl.compile(r'^\[\d{2}:\d{2}\]')
_NAME_KW = _re_pl.compile(
    r'(transcript|trascrizione|lezione|audio|registr|carne)',
    _re_pl.IGNORECASE
)


def _strip_rtf(raw: str) -> str:
    """Rimuove markup RTF e restituisce testo pulito."""
    if not raw.lstrip().startswith('{\\rtf'):
        return raw
    try:
        from striprtf.striprtf import rtf_to_text
        return rtf_to_text(raw)
    except ImportError:
        pass
    t = _re_pl.sub(r'\\\n', '\n', raw)
    t = _re_pl.sub(r'\\par\b', '\n', t)
    t = _re_pl.sub(r'\\line\b', '\n', t)
    t = _re_pl.sub(r'\\\\\s?', '', t)
    t = _re_pl.sub(r'\\[a-z]+\-?\d*\s?', '', t)
    t = _re_pl.sub(r'\{[^{}]{0,200}\}', '', t)
    t = _re_pl.sub(r'[{}]', '', t)
    t = _re_pl.sub(r"\\\'([0-9a-f]{2})",
                   lambda m: bytes.fromhex(m.group(1)).decode('cp1252', errors='replace'), t)
    t = _re_pl.sub(r'\n{3,}', '\n\n', t)
    return t.strip()


def _collect_audio(sources: dict, source_names: list, audio_files: list,
                   output_dir: Path, whisper_model: str) -> None:
    """STEP 1: Trascrive file audio → CARNE."""
    for idx, af in enumerate(audio_files):
        print(f"\n  [Audio → CARNE] {af.name}")
        prog_base = 10 + idx * 5
        _report_progress(output_dir, prog_base,
                         "Whisper — Trascrizione audio",
                         f"File {idx+1}/{len(audio_files)}: {af.name}")
        t, dur = transcribe_audio(af, whisper_model,
                                  _prog_dir=output_dir,
                                  _prog_base=prog_base, _prog_end=prog_base + 5)
        if t:
            sources["carne"].append({"filename": af.name, "text": t, "duration_sec": dur})
            sources["has_audio"] = True
            source_names.append(af.name)


def _collect_video(sources: dict, source_names: list, video_files: list,
                   output_dir: Path, tmp_dir: Path, whisper_model: str) -> None:
    """STEP 2: Estrae audio da video e trascrive → CARNE."""
    for idx, vf in enumerate(video_files):
        print(f"\n  [Video → CARNE] {vf.name}")
        prog_base = 15 + idx * 5
        _report_progress(output_dir, prog_base,
                         "Whisper — Trascrizione video",
                         f"File {idx+1}/{len(video_files)}: {vf.name}")
        mp3 = extract_audio_from_video(vf, tmp_dir)
        if mp3:
            t, dur = transcribe_audio(mp3, whisper_model,
                                      _prog_dir=output_dir,
                                      _prog_base=prog_base, _prog_end=prog_base + 5)
            if t:
                sources["carne"].append({"filename": vf.name, "text": t, "duration_sec": dur})
                sources["has_audio"] = True
                source_names.append(vf.name)


def _collect_pptx(sources: dict, source_names: list, slide_files: list,
                  images_dir: Path, lesson_number: int, skip_ocr: bool):
    """STEP 3: Estrae slide PPTX → SCHELETRO. Ritorna l'oggetto pptx_slides (o None)."""
    pptx_slides = None
    for sf in slide_files:
        print(f"\n  [PPTX → SCHELETRO] {sf.name}")
        _report_progress(images_dir.parent, 35, "Estrazione slide PPTX", sf.name)
        if COLLEAGUE_MODULES:
            slides_obj, plain = process_pptx_full(sf, images_dir, skip_ocr=skip_ocr)
            pptx_slides = slides_obj
            if len(slide_files) > 1:
                pptx_img_dir = images_dir / sf.stem
                pptx_img_dir.mkdir(exist_ok=True)
                slide_images_raw = render_slide_images(sf, pptx_img_dir)
                slide_images = {k: f"{sf.stem}/{v}" for k, v in slide_images_raw.items()}
            else:
                slide_images = render_slide_images(sf, images_dir)
            skeleton_latex = build_fallback_latex(
                lesson_number = lesson_number,
                title         = sf.stem.replace("_", " ").replace("-", " ").title(),
                slides        = slides_obj,
                transcript    = "",
                slide_text    = plain,
                extra_text    = "",
                slide_images  = slide_images,
            )
            sources["scheletro"].append({
                "filename":    sf.name,
                "text":        plain,
                "latex":       skeleton_latex,
                "slide_count": len(slides_obj),
                "slide_images": slide_images,
            })
        else:
            plain = process_pptx_fallback(sf)
            sources["scheletro"].append({
                "filename":    sf.name,
                "text":        plain,
                "latex":       None,
                "slide_count": plain.count("--- SLIDE"),
            })
        source_names.append(sf.name)
    return pptx_slides


def _collect_txt(sources: dict, source_names: list, text_files: list,
                 has_structure: bool, has_real_audio: bool) -> None:
    """STEP 4: Classifica file TXT/MD come CARNE (trascrizioni) o CONTORNO (note)."""
    for tf in text_files:
        raw  = tf.read_text(encoding="utf-8", errors="ignore")
        text = _strip_rtf(raw)
        if not text.strip():
            continue
        lines = [l for l in text.splitlines() if l.strip()]
        ts_count = sum(1 for l in lines if _TS_PAT.match(l.strip()))
        if len(lines) == 0:
            has_timestamps = False
        elif len(lines) <= 10:
            has_timestamps = ts_count / len(lines) >= 0.30
        else:
            has_timestamps = ts_count / len(lines) >= 0.10
        has_kw_name      = bool(_NAME_KW.search(tf.stem))
        is_solo_companion = has_structure and not has_real_audio and not sources["has_audio"]
        is_transcript     = has_timestamps or has_kw_name or is_solo_companion
        if is_transcript:
            reason = ("timestamp" if has_timestamps
                      else "nome file" if has_kw_name
                      else "unico txt + struttura presente")
            print(f"\n  [TXT → CARNE] {tf.name}  ({reason})")
            sources["carne"].append({"filename": tf.name, "text": text, "duration_sec": 0})
            sources["has_audio"] = True
        else:
            print(f"\n  [TXT → CONTORNO] {tf.name}")
            sources["contorno"].append({"filename": tf.name, "text": text})
        source_names.append(tf.name)


def _collect_pdf(sources: dict, source_names: list, pdf_files: list,
                 images_dir: Path, output_dir: Path, lesson_number: int,
                 title: str, skip_ai: bool, skip_ocr: bool,
                 subject_hint: str, course_context_path: str):
    """STEP 5: Estrae PDF → SCHELETRO/SUPPORTO, o chunking se >20 pag senza audio.
    Ritorna list[Path] se chunking (early exit), None altrimenti."""
    _PDF_CHUNK_THRESHOLD = 20
    for idx, pf in enumerate(pdf_files):
        _report_progress(output_dir, 40 + idx * 3, "Estrazione pagine PDF", pf.name)
        pages = extract_pdf_pages(pf)
        if COLLEAGUE_MODULES:
            page_images, latex_skeleton = render_pdf_pages(pf, images_dir, pages or None)
        else:
            page_images, latex_skeleton = {}, None
        if page_images:
            _pages_with_text = {p["page"] for p in pages}
            _missing = {n: f for n, f in page_images.items() if n not in _pages_with_text}
            if _missing:
                _report_progress(output_dir, 40 + idx * 3, "OCR pagine scansionate", pf.name)
                _ocr = _ocr_pages_with_tesseract(images_dir, _missing)
                if _ocr:
                    pages.extend(_ocr)
                    pages.sort(key=lambda p: p["page"])
                    if COLLEAGUE_MODULES:
                        latex_skeleton = _build_pdf_latex_skeleton(pf, pages, page_images)
        if not pages:
            print(f"    [WARN] {pf.name}: nessun testo disponibile — saltato")
            continue
        if len(pages) > _PDF_CHUNK_THRESHOLD and not sources["has_audio"]:
            print(f"\n  [PDF grande] {pf.name} — {len(pages)} pagine, chunking automatico")
            return process_pdf_chunked(
                pdf_path            = pf,
                output_dir          = output_dir,
                base_lesson_number  = lesson_number,
                title               = title or pf.stem,
                skip_ai             = skip_ai,
                subject_hint        = subject_hint,
                course_context_path = course_context_path,
            )
        text  = "\n".join(f"[PAG {p['page']}]\n{p['text']}" for p in pages)
        entry = {
            "filename":    pf.name,
            "text":        text,
            "pages":       len(pages),
            "latex":       latex_skeleton,
            "page_images": page_images,
        }
        if sources["has_audio"]:
            print(f"\n  [PDF → SCHELETRO] {pf.name}  (c'è audio)")
            sources["scheletro"].append(entry)
        else:
            print(f"\n  [PDF → SUPPORTO] {pf.name}  (nessun audio)")
            sources["supporto"].append(entry)
        source_names.append(pf.name)
    return None


def _collect_docx(sources: dict, source_names: list, doc_files: list,
                  output_dir: Path) -> None:
    """STEP 6: Estrae DOCX → SCHELETRO se c'è audio, SUPPORTO altrimenti."""
    for idx, df in enumerate(doc_files):
        _report_progress(output_dir, 50 + idx * 3, "Estrazione DOCX", df.name)
        text = extract_docx(df)
        if not text:
            continue
        entry = {"filename": df.name, "text": text, "latex": None, "pages": None}
        if sources["has_audio"]:
            print(f"\n  [DOCX → SCHELETRO] {df.name}  (c'è audio)")
            sources["scheletro"].append(entry)
        else:
            print(f"\n  [DOCX → SUPPORTO] {df.name}  (nessun audio)")
            sources["supporto"].append(entry)
        source_names.append(df.name)


# ─────────────────────────────────────────────
# CORE: PROCESSA UNA LEZIONE
# ─────────────────────────────────────────────

def process_lesson(source_dir: Path, lesson_number: int, output_dir: Path,
                   skip_ai: bool = False, skip_ocr: bool = False,
                   whisper_model: str = "base",
                   subject_hint: str = None,
                   course_context_path: str = None,
                   title: str = None):

    print(f"\n{'─'*58}")
    print(f"  LEZIONE {lesson_number:02d}  ←  {source_dir.name}")
    print(f"{'─'*58}")

    # Raccogli file: escludi file nascosti/sistema e ordina per nome (output deterministico)
    _SKIP_NAMES = {".ds_store", "thumbs.db", "desktop.ini", ".gitkeep", ".gitignore"}
    all_files = sorted(
        (
            f for f in (source_dir.iterdir() if source_dir.is_dir() else [source_dir])
            if f.is_file()
            and not f.name.startswith(".")
            and not f.name.startswith("__")
            and f.name.lower() not in _SKIP_NAMES
        ),
        key=lambda p: p.name.lower(),
    )

    by = {k: [] for k in ("audio", "video", "slide", "doc", "pdf", "text")}
    for f in all_files:
        ext = f.suffix.lower()
        if   ext in CONFIG["ext_audio"]: by["audio"].append(f)
        elif ext in CONFIG["ext_video"]: by["video"].append(f)
        elif ext in CONFIG["ext_slide"]: by["slide"].append(f)
        elif ext in CONFIG["ext_doc"]:   by["doc"].append(f)
        elif ext in CONFIG["ext_pdf"]:   by["pdf"].append(f)
        elif ext in CONFIG["ext_text"]:  by["text"].append(f)

    labels = (
        [f"audio:{f.name}" for f in by["audio"]] +
        [f"video:{f.name}" for f in by["video"]] +
        [f"pptx:{f.name}"  for f in by["slide"]] +
        [f"docx:{f.name}"  for f in by["doc"]]   +
        [f"pdf:{f.name}"   for f in by["pdf"]]   +
        [f"txt:{f.name}"   for f in by["text"]]
    )
    if not labels:
        print("  [SKIP] Nessun file riconosciuto")
        return None
    print(f"  Fonti trovate: {', '.join(labels)}")
    _report_progress(output_dir, 5, "Analisi sorgenti", f"Lezione {lesson_number}")

    tmp_dir    = output_dir / f"_tmp_{lesson_number:02d}"
    images_dir = output_dir / CONFIG["images_subdir"]
    tmp_dir.mkdir(exist_ok=True)
    images_dir.mkdir(exist_ok=True)

    sources = {
        "scheletro": [],
        "carne":     [],
        "supporto":  [],
        "contorno":  [],
        "has_audio": False,
    }
    source_names  = []

    _collect_audio(sources, source_names, by["audio"], output_dir, whisper_model)
    _collect_video(sources, source_names, by["video"], output_dir, tmp_dir, whisper_model)
    pptx_slides = _collect_pptx(sources, source_names, by["slide"], images_dir, lesson_number, skip_ocr)
    _collect_txt(sources, source_names, by["text"],
                 has_structure=bool(by["slide"] or by["pdf"] or by["doc"]),
                 has_real_audio=bool(by["audio"] or by["video"]))
    early_result = _collect_pdf(sources, source_names, by["pdf"], images_dir, output_dir,
                                lesson_number, title, skip_ai, skip_ocr,
                                subject_hint, course_context_path)
    if early_result is not None:
        return early_result
    _collect_docx(sources, source_names, by["doc"], output_dir)

    # ── Gerarchia risolta ──
    print(f"\n  Gerarchia risolta:")
    print(f"    SCHELETRO : {[s['filename'] for s in sources['scheletro']] or '—'}")
    print(f"    CARNE     : {[s['filename'] for s in sources['carne']] or '—'}")
    print(f"    SUPPORTO  : {[s['filename'] for s in sources['supporto']] or '—'}")
    print(f"    CONTORNO  : {[s['filename'] for s in sources['contorno']] or '—'}")

    # Salva riepilogo classificazione nel debug dir
    try:
        _dbg_dir = output_dir / "debug"
        _dbg_dir.mkdir(parents=True, exist_ok=True)
        _dbg_report_path = _dbg_dir / f"riepilogo_lezione_{lesson_number:02d}.txt"
        _char_count = lambda lst: sum(len(s.get("text","")) for s in lst)
        _report_lines = [
            f"Lezione {lesson_number:02d} — {source_dir.name}",
            f"Processata: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
            "",
            "FILE TROVATI:",
        ]
        for f in all_files:
            _report_lines.append(f"  {f.name}  ({f.stat().st_size // 1024 + 1} KB)")
        _report_lines += [
            "",
            "CLASSIFICAZIONE SEMANTICA:",
            f"  SCHELETRO ({len(sources['scheletro'])} file, ~{_char_count(sources['scheletro']):,} char):",
        ]
        for s in sources["scheletro"]:
            _report_lines.append(f"    • {s['filename']}")
        _report_lines.append(f"  CARNE ({len(sources['carne'])} file, ~{_char_count(sources['carne']):,} char):")
        for s in sources["carne"]:
            _report_lines.append(f"    • {s['filename']}  (dur: {s.get('duration_sec',0)}s)")
        _report_lines.append(f"  SUPPORTO ({len(sources['supporto'])} file, ~{_char_count(sources['supporto']):,} char):")
        for s in sources["supporto"]:
            _report_lines.append(f"    • {s['filename']}")
        _report_lines.append(f"  CONTORNO ({len(sources['contorno'])} file, ~{_char_count(sources['contorno']):,} char):")
        for s in sources["contorno"]:
            _report_lines.append(f"    • {s['filename']}")
        _dbg_report_path.write_text("\n".join(_report_lines), encoding="utf-8")
    except Exception:
        pass

    _resize_images_dir(images_dir)

    if not title:
        title = source_dir.name.replace("_", " ").replace("-", " ").title()
    out_tex = output_dir / f"lezione_{lesson_number:02d}.tex"

    def _latex_from_skeleton(sources, lesson_number, title, pptx_slides):
        entry = (
            next((s for s in sources["scheletro"] if s.get("latex")), None) or
            next((s for s in sources["supporto"]  if s.get("latex")), None)
        )
        if entry:
            return (
                f"\\section{{Lezione {lesson_number}: {title}}}\n"
                f"\\label{{sec:lezione{lesson_number:02d}}}\n\n"
                + entry["latex"]
            )
        slide_text = "\n".join(s["text"] for s in sources["scheletro"])
        carne_text = "\n".join(s["text"] for s in sources["carne"])
        extra_text = "\n".join(s["text"] for s in sources["supporto"] + sources["contorno"])
        return build_fallback_latex(
            lesson_number, title, pptx_slides,
            carne_text, slide_text, extra_text,
        )

    print()
    _report_progress(output_dir, 60, "Claude — Costruzione prompt", "")
    content = generate_with_claude(
        lesson_number        = lesson_number,
        title                = title,
        sources              = sources,
        subject_hint         = subject_hint,
        course_context_path  = course_context_path,
        _progress_output_dir = output_dir,
        skip_ai              = skip_ai,
    )
    if not content:
        if skip_ai:
            print("  [--skip-ai] LaTeX strutturato senza Claude")
        else:
            print("  [fallback] Claude non raggiunto, uso scheletro")
        _report_progress(output_dir, 80, "Generazione LaTeX (fallback scheletro)", "")
        content = _latex_from_skeleton(sources, lesson_number, title, pptx_slides)

    _report_progress(output_dir, 90, "Scrittura file LaTeX", out_tex.name)
    write_lesson_tex(lesson_number, title, content, source_names, out_tex)
    shutil.rmtree(tmp_dir, ignore_errors=True)
    return out_tex
# ─────────────────────────────────────────────
# ENTRY POINT
# ─────────────────────────────────────────────
def main():
    parser = argparse.ArgumentParser(
        description="Appunti AI — fonti miste -> libro LaTeX",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__
    )
    parser.add_argument("source",
        help="Cartella lezione (singola) o cartella corso (con --batch)")
    parser.add_argument("--batch", action="store_true",
        help="Ogni sottocartella = una lezione")
    parser.add_argument("--output", default="./output",
        help="Cartella di output (default: ./output)")
    parser.add_argument("--title", default="Appunti del Corso",
        help="Titolo del corso per main.tex")
    parser.add_argument("--skip-ai", action="store_true",
        help="Non usare Claude (usa struttura dal collega, offline)")
    parser.add_argument("--skip-ocr", action="store_true",
        help="Non usare pix2tex OCR (più veloce)")
    parser.add_argument("--whisper-model", default=CONFIG["whisper_model"],
        choices=["tiny","base","small","medium","large"],
        help="Modello Whisper (default: base)")
    parser.add_argument("--start-from", type=int, default=None,
        help="Forza la numerazione a partire da N (default: auto da state.json)")
    parser.add_argument("--subject",
        choices=["ingegneria","matematica","fisica","medicina",
                "economia","giurisprudenza","generico"],
        help="Tipo di materia (default: auto-detect)")
    parser.add_argument("--no-context", action="store_true",
        help="Non usare/aggiornare corso_context.json")
    parser.add_argument("--continue-on-error", action="store_true",
        help="In batch: salta lezioni che falliscono invece di interrompere tutto")
    args = parser.parse_args()

    source_path  = Path(args.source)
    output_dir   = Path(args.output)
    output_dir.mkdir(parents=True, exist_ok=True)
    (output_dir / CONFIG["images_subdir"]).mkdir(exist_ok=True)

    # Path contesto corso (auto, nella cartella output)
    course_context_path = None
    if PREPROCESSOR and not args.no_context and not args.skip_ai:
        course_context_path = str(output_dir / "corso_context.json")

    print(f"\n{'═'*58}")
    print(f"  APPUNTI AI")
    print(f"  Output  : {output_dir.resolve()}")
    print(f"  Titolo  : {args.title}")
    ai_str  = "OFF (--skip-ai)" if args.skip_ai else f"Claude {CONFIG['claude_model']}"
    ocr_str = "OFF (--skip-ocr)" if args.skip_ocr else "pix2tex"
    sub_str = args.subject or "auto-detect"
    ctx_str = course_context_path or "OFF (--no-context)"
    print(f"  Claude  : {ai_str}")
    print(f"  Whisper : {args.whisper_model}")
    print(f"  OCR     : {ocr_str}")
    print(f"  Materia : {sub_str}")
    print(f"  Contesto: {ctx_str}")
    print(f"{'═'*58}")

    # ── Carica stato persistente del corso ──────────────────
    state = load_state(output_dir)
    if args.title and args.title != "Appunti del Corso":
        state["course_title"] = args.title
    if args.subject:
        state["subject"] = args.subject

    # Determina da quale numero partire
    if args.start_from is not None:
        # Override esplicito → rispetta sempre (utile per correzioni)
        next_lesson = args.start_from
        print(f"  Numerazione : forzata da {next_lesson} (--start-from)")
    else:
        next_lesson = state["next_lesson"]
        print(f"  Numerazione : auto da state.json (prossima lezione {next_lesson})")

    lesson_files = []

    def collect(result, source_dir_name: str = ""):
        """Gestisce sia Path singolo che lista di Path (PDF chunked).
        Aggiorna state.json dopo ogni lezione processata con successo."""
        if result is None:
            return
        results = result if isinstance(result, list) else [result]
        for tex_path in results:
            lesson_files.append(tex_path)
            # Aggiorna lista lezioni nello stato
            # Evita duplicati: rimuovi eventuale voce precedente con stesso tex_file
            state["lessons"] = [
                e for e in state["lessons"] if e["tex_file"] != tex_path.name
            ]
            state["lessons"].append({
                "number":       int(tex_path.stem.split("_")[-1]),
                "source_dir":   source_dir_name,
                "tex_file":     tex_path.name,
                "processed_at": datetime.now().isoformat(timespec="seconds"),
            })
        # Ricalcola next_lesson come max numero usato + 1
        if state["lessons"]:
            state["next_lesson"] = max(e["number"] for e in state["lessons"]) + 1
        save_state(output_dir, state)
        print(f"  [state] salvato → prossima lezione: {state['next_lesson']}")

    batch_errors: list[dict] = []

    if args.batch:
        subdirs = sorted(d for d in source_path.iterdir() if d.is_dir())
        if not subdirs:
            print(f"[ERRORE] Nessuna sottocartella in {source_path}")
            sys.exit(1)
        print(f"\n  {len(subdirs)} lezioni trovate")
        processed_dirs = {e["source_dir"] for e in state["lessons"]}
        for subdir in subdirs:
            if subdir.name in processed_dirs:
                print(f"  [SKIP] {subdir.name} già processata (presente in state.json)")
                continue
            lesson_num = state["next_lesson"]
            try:
                result = process_lesson(
                    subdir, lesson_num, output_dir,
                    skip_ai             = args.skip_ai,
                    skip_ocr            = args.skip_ocr,
                    whisper_model       = args.whisper_model,
                    subject_hint        = args.subject,
                    course_context_path = course_context_path,
                    # In batch ogni sottocartella ha nome semantico → fallback al nome cartella
                    title               = None,
                )
                collect(result, subdir.name)
            except Exception as exc:
                if args.continue_on_error:
                    err_info = {"subdir": subdir.name, "error": str(exc)}
                    batch_errors.append(err_info)
                    print(f"\n  [SKIP-ERROR] {subdir.name}: {exc}")
                    continue
                raise
    else:
        result = process_lesson(
            source_path, next_lesson, output_dir,
            skip_ai             = args.skip_ai,
            skip_ocr            = args.skip_ocr,
            whisper_model       = args.whisper_model,
            subject_hint        = args.subject,
            course_context_path = course_context_path,
            title               = args.title if args.title != "Appunti del Corso" else None,
        )
        collect(result, source_path.name)

    # ── main.tex: include TUTTE le lezioni dello stato, non solo quella corrente
    all_tex_files = []
    for entry in sorted(state["lessons"], key=lambda e: e["number"]):
        tex_path = output_dir / entry["tex_file"]
        if tex_path.exists():
            all_tex_files.append(tex_path)

    if all_tex_files:
        _report_progress(output_dir, 95, "Generazione main.tex",
                         f"{len(all_tex_files)} lezioni")
        generate_main_tex(
            state["course_title"] or args.title,
            all_tex_files,
            output_dir,
            lang = os.environ.get("WHISPER_LANG") or state.get("subject_lang") or "en",
        )
        # Salva errori batch se --continue-on-error era attivo
        if batch_errors:
            import json as _json
            err_file = output_dir / "errors.json"
            err_file.write_text(
                _json.dumps(batch_errors, indent=2, ensure_ascii=False),
                encoding="utf-8"
            )
            print(f"\n  [ATTENZIONE] {len(batch_errors)} lezioni saltate → errors.json")
        _report_progress(output_dir, 100, "Completato", "")
        print(f"\n{'═'*58}")
        print(f"  COMPLETATO — {len(lesson_files)} lezioni"
              + (f" ({len(batch_errors)} errori)" if batch_errors else ""))
        print(f"\n  Compila il PDF:")
        print(f"    cd {output_dir.resolve()}")
        print(f"    pdflatex main.tex && pdflatex main.tex")
        print(f"{'═'*58}\n")
    else:
        print("\n[ATTENZIONE] Nessun .tex generato.")
        sys.exit(1)


if __name__ == "__main__":
    main()