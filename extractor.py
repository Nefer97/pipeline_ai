"""
extractor.py
Estrae testo, immagini e formule OMML da un file .pptx
"""

import os
import re
import subprocess
from dataclasses import dataclass, field
from typing import Optional
from pptx import Presentation
from pptx.util import Emu
from lxml import etree
import hashlib

# Formati non supportati da LaTeX/Pillow → da convertire in PNG
UNSUPPORTED_FORMATS = {'.wmf', '.emf', '.gif', '.bmp', '.tiff', '.tif'}


def _convert_to_png(src_path: str) -> str:
    """
    Converte un'immagine in PNG.
    wmf/emf: usa LibreOffice o Inkscape.
    gif/bmp/tiff: usa Pillow.
    Ritorna il path del PNG, o src_path originale se la conversione fallisce.
    """
    ext = os.path.splitext(src_path)[1].lower()
    abs_src  = os.path.abspath(src_path)
    png_path = abs_src.rsplit('.', 1)[0] + '.png'

    if ext in ('.wmf', '.emf'):
        # Prova LibreOffice
        # os.path.dirname può restituire stringa vuota per path relativi → usa abspath
        out_dir = os.path.dirname(os.path.abspath(src_path))
        try:
            result = subprocess.run(
                ['libreoffice', '--headless', '--convert-to', 'png',
                 '--outdir', out_dir, src_path],
                capture_output=True, timeout=30
            )
            if result.returncode == 0 and os.path.exists(png_path):
                try:
                    os.remove(src_path)
                except OSError:
                    pass  # file in uso o già rimosso — non critico
                return png_path
        except (subprocess.TimeoutExpired, FileNotFoundError):
            pass
        # Prova Inkscape (png_path è già assoluto grazie all'abspath sopra)
        try:
            result = subprocess.run(
                ['inkscape', '--export-type=png', f'--export-filename={png_path}', src_path],
                capture_output=True, timeout=30
            )
            if result.returncode == 0 and os.path.exists(png_path):
                try:
                    os.remove(src_path)
                except OSError:
                    pass
                return png_path
        except (subprocess.TimeoutExpired, FileNotFoundError):
            pass
        print(f"  [WARN] {os.path.basename(src_path)}: wmf/emf non convertibile (installa LibreOffice o Inkscape)")
        return src_path
    else:
        # gif, bmp, tiff → Pillow
        try:
            from PIL import Image
            img = Image.open(src_path)
            img.convert('RGBA').save(png_path, 'PNG')
            try:
                os.remove(src_path)
            except OSError:
                pass
            return png_path
        except Exception as e:
            print(f"  [WARN] Conversione {os.path.basename(src_path)} fallita: {e}")
            return src_path


# Namespace XML usati in pptx
NSMAP = {
    'a':   'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p':   'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r':   'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'm':   'http://schemas.openxmlformats.org/officeDocument/2006/math',
    'pic': 'http://schemas.openxmlformats.org/drawingml/2006/picture',
}


@dataclass
class SlideObject:
    """Un oggetto estratto da una slide, con posizione."""
    obj_type: str          # 'text', 'image', 'omml_formula'
    content: str           # testo o path immagine o OMML xml string
    top: float             # coordinata Y in EMU
    left: float            # coordinata X in EMU
    width: float = 0.0
    height: float = 0.0
    image_path: Optional[str] = None  # solo per immagini


@dataclass
class SlideData:
    slide_number: int
    title: str
    objects: list = field(default_factory=list)  # lista di SlideObject ordinati per Y
    notes: str = ""                              # note del presentatore (testo dalla Notes pane)


def _emu_to_pt(emu: int) -> float:
    return emu / 12700.0


def _get_position(shape):
    """Restituisce (top, left, width, height) in EMU."""
    try:
        return shape.top or 0, shape.left or 0, shape.width or 0, shape.height or 0
    except Exception:
        return 0, 0, 0, 0


def _extract_omml(shape_xml: etree._Element) -> Optional[str]:
    """Cerca formula OMML (m:oMath) nell'elemento XML della shape."""
    math_elems = shape_xml.findall('.//' + '{http://schemas.openxmlformats.org/officeDocument/2006/math}oMath')
    if math_elems:
        return etree.tostring(math_elems[0], encoding='unicode')
    return None


def _extract_text(shape) -> str:
    """Estrae testo pulito da una shape testuale."""
    try:
        if shape.has_text_frame:
            lines = []
            for para in shape.text_frame.paragraphs:
                line = ''.join(run.text for run in para.runs)
                if line.strip():
                    lines.append(line)
            return '\n'.join(lines)
    except Exception:
        pass
    return ''


def _extract_table_latex(table) -> str:
    """
    Converte una tabella python-pptx in LaTeX \\begin{tabular}.
    Ritorna stringa LaTeX o '' se la tabella è vuota/non leggibile.
    """
    try:
        rows_data = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                text = cell.text_frame.text.strip().replace("\n", " ") if cell.text_frame else ""
                # Escape caratteri speciali LaTeX nel contenuto della cella
                for old, new in [("&", "\\&"), ("%", "\\%"), ("$", "\\$"),
                                  ("#", "\\#"), ("_", "\\_"), ("{", "\\{"), ("}", "\\}")]:
                    text = text.replace(old, new)
                cells.append(text)
            rows_data.append(cells)

        if not rows_data:
            return ""

        n_cols = max(len(r) for r in rows_data)
        if n_cols == 0:
            return ""

        col_spec = "|" + "l|" * n_cols

        lines = [
            "\\begin{center}",
            f"\\begin{{tabular}}{{{col_spec}}}",
            "\\hline",
        ]
        for i, row in enumerate(rows_data):
            # Padding celle mancanti
            padded = row + [""] * (n_cols - len(row))
            lines.append(" & ".join(padded) + " \\\\")
            # Riga separatrice dopo intestazione (prima riga) e alla fine
            if i == 0 or i == len(rows_data) - 1:
                lines.append("\\hline")
        lines += ["\\end{tabular}", "\\end{center}"]
        return "\n".join(lines)
    except Exception:
        return ""


def extract_slides(pptx_path: str, image_output_dir: str) -> list:
    """
    Estrae tutti gli oggetti da ogni slide.
    Ritorna lista di SlideData.
    """
    os.makedirs(image_output_dir, exist_ok=True)
    prs = Presentation(pptx_path)
    slides_data = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        # Titolo
        title = ''
        try:
            title = (slide.shapes.title.text or '') if slide.shapes.title else ''
        except Exception:
            pass

        objects = []

        for shape in slide.shapes:
            top, left, width, height = _get_position(shape)

            # --- Controlla prima se è formula OMML ---
            omml = _extract_omml(shape._element)
            if omml:
                objects.append(SlideObject(
                    obj_type='omml_formula',
                    content=omml,
                    top=top, left=left, width=width, height=height
                ))
                continue

            # --- Immagine ---
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    image = shape.image
                    ext = image.ext  # png, jpg, ...
                    # Nome file basato su hash per evitare duplicati
                    img_hash = hashlib.md5(image.blob).hexdigest()[:8]
                    img_filename = f"slide{slide_idx:03d}_{img_hash}.{ext}"
                    img_path = os.path.join(image_output_dir, img_filename)
                    with open(img_path, 'wb') as f:
                        f.write(image.blob)
                    # Converti formati non supportati da LaTeX.
                    # image.ext restituisce l'estensione SENZA punto (es. "gif"),
                    # mentre UNSUPPORTED_FORMATS usa punti (es. ".gif").
                    if ('.' + ext.lower()) in UNSUPPORTED_FORMATS:
                        img_path = _convert_to_png(img_path)
                        img_filename = os.path.basename(img_path)
                    objects.append(SlideObject(
                        obj_type='image',
                        content=img_filename,
                        top=top, left=left, width=width, height=height,
                        image_path=img_path
                    ))
                except Exception as e:
                    print(f"  [WARN] Immagine non estratta slide {slide_idx}: {e}")
                continue

            # --- Tabella ---
            try:
                from pptx.enum.shapes import PP_PLACEHOLDER
                if shape.has_table:
                    table_latex = _extract_table_latex(shape.table)
                    if table_latex:
                        objects.append(SlideObject(
                            obj_type='table',
                            content=table_latex,
                            top=top, left=left, width=width, height=height
                        ))
                    continue
            except Exception:
                pass

            # --- Testo ---
            text = _extract_text(shape)
            if text.strip():
                objects.append(SlideObject(
                    obj_type='text',
                    content=text,
                    top=top, left=left, width=width, height=height
                ))

        # Ordina per posizione verticale (Y crescente)
        objects.sort(key=lambda o: o.top)

        # Note del presentatore (Notes pane)
        notes = ""
        try:
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                if notes_tf:
                    # Filtra placeholder PowerPoint in varie lingue e simboli decorativi
                    raw = notes_tf.text.strip()
                    _placeholder_prefixes = (
                        "click to edit", "fare clic per modificare",
                        "haga clic para editar", "zum bearbeiten klicken",
                        "cliquez pour modifier",
                    )
                    _stripped = raw.replace("*","").replace("—","").replace("-","").replace(".","").strip()
                    _is_placeholder = (
                        not raw
                        or raw.lower().startswith(_placeholder_prefixes)
                        or _stripped == ""
                        or len(raw) < 3
                    )
                    if not _is_placeholder:
                        notes = raw
        except Exception:
            pass

        slides_data.append(SlideData(
            slide_number=slide_idx,
            title=title,
            objects=objects,
            notes=notes,
        ))

    return slides_data
